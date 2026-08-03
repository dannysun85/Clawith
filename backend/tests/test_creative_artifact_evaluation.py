from __future__ import annotations

from pathlib import Path
from types import SimpleNamespace

import fitz
from PIL import Image
import pytest
from pptx import Presentation
from pptx.util import Inches

from app.services.creative_artifact_evaluation import (
    CreativeArtifactContract,
    observe_creative_artifacts,
    score_artifact_observation,
)
from app.services.creative_evaluation import generate_evaluation_bundle


def _write_image(path: Path, *, size: tuple[int, int] = (900, 1600)) -> None:
    image = Image.new("RGB", size, color=(32, 64, 96))
    image.save(path)


def _write_pdf(path: Path, *, page_count: int = 2) -> None:
    document = fitz.open()
    for _ in range(page_count):
        document.new_page(width=960, height=540)
    document.save(path)
    document.close()


def _write_editable_pptx(path: Path, *, page_count: int = 2) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    for index in range(page_count):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(8),
            Inches(1),
        )
        box.text = f"Slide {index + 1}"
    presentation.save(path)


def _write_screenshot_pptx(path: Path, image_path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        presentation.slide_width,
        presentation.slide_height,
    )
    presentation.save(path)


def _write_picture_mix_pptx(
    path: Path,
    image_path: Path,
    *,
    full_slide_first: bool,
) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    for index in range(2):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        if index == 0 and full_slide_first:
            slide.shapes.add_picture(
                str(image_path),
                0,
                0,
                presentation.slide_width,
                presentation.slide_height,
            )
        elif index == 0:
            slide.shapes.add_picture(
                str(image_path),
                Inches(0.5),
                Inches(0.5),
                Inches(2),
                Inches(1),
            )
        box = slide.shapes.add_textbox(
            Inches(1),
            Inches(1.5),
            Inches(8),
            Inches(1),
        )
        box.text = f"Slide {index + 1}"
    presentation.save(path)


@pytest.mark.asyncio
async def test_image_observation_proves_structure_but_not_semantics(tmp_path) -> None:
    path = tmp_path / "candidate.png"
    _write_image(path)

    observation = await observe_creative_artifacts(
        CreativeArtifactContract(modality="image", aspect_ratio="9:16"),
        {"image": path},
    )

    assert observation.hard_gates["artifact_decodable"].passed is True
    assert observation.hard_gates["aspect_ratio_match"].passed is True
    assert observation.hard_gates["fact_safety"].passed is None
    assert observation.hard_gates["no_unrequested_watermark"].passed is None
    assert observation.facts["width"] == 900
    assert observation.facts["height"] == 1600


@pytest.mark.asyncio
async def test_image_wrong_aspect_ratio_fails_hard_gate(tmp_path) -> None:
    path = tmp_path / "candidate.png"
    _write_image(path, size=(1200, 1200))

    observation = await observe_creative_artifacts(
        CreativeArtifactContract(modality="image", aspect_ratio="9:16"),
        {"image": path},
    )

    assert observation.hard_gates["aspect_ratio_match"].passed is False

    scenario = next(
        item
        for item in generate_evaluation_bundle(seed=20260727).manifest.public_scenarios
        if item.modality == "image"
    )
    scenario = scenario.model_copy(update={"aspect_ratio": "9:16"})
    evaluation = score_artifact_observation(scenario, observation)
    assert evaluation.status == "blocked"
    assert "aspect_ratio_match" in evaluation.hard_gate_failures


@pytest.mark.asyncio
async def test_video_observation_checks_delivery_contract(
    tmp_path,
    monkeypatch,
) -> None:
    path = tmp_path / "candidate.mp4"
    path.write_bytes(b"\x00\x00\x00\x18ftypmp42video")
    info = SimpleNamespace(
        width=720,
        height=1280,
        duration_seconds=6.0,
        codec_name="h264",
        pixel_format="yuv420p",
        audio_codec_name="aac",
        fast_start=True,
    )

    async def validate_video(*_args, **_kwargs):
        return info

    monkeypatch.setattr(
        "app.services.media_assets.validate_generated_video",
        validate_video,
    )
    monkeypatch.setattr(
        "app.services.media_assets.validate_video_delivery_contract",
        lambda value, **_kwargs: value,
    )

    observation = await observe_creative_artifacts(
        CreativeArtifactContract(
            modality="video",
            aspect_ratio="9:16",
            duration_seconds=6,
            audio_required=True,
        ),
        {"mp4": path},
    )

    assert observation.hard_gates["artifact_decodable"].passed is True
    assert observation.hard_gates["duration_and_aspect_match"].passed is True
    assert observation.hard_gates["audio_contract_match"].passed is True
    assert observation.hard_gates["fact_safety"].passed is None


@pytest.mark.asyncio
async def test_presentation_observation_accepts_editable_pptx_and_pdf(
    tmp_path,
) -> None:
    pptx_path = tmp_path / "deck.pptx"
    pdf_path = tmp_path / "deck.pdf"
    _write_editable_pptx(pptx_path)
    _write_pdf(pdf_path)

    observation = await observe_creative_artifacts(
        CreativeArtifactContract(
            modality="presentation",
            aspect_ratio="16:9",
            page_count=2,
            editable_required=True,
        ),
        {"pptx": pptx_path, "pdf": pdf_path},
    )

    assert observation.hard_gates["pptx_and_preview_valid"].passed is True
    assert observation.hard_gates["page_count_and_aspect_match"].passed is True
    assert observation.hard_gates["editability"].passed is True
    assert observation.hard_gates["no_text_overflow"].passed is None
    assert observation.facts["pptx"]["editable_slide_ratio"] == 1
    assert observation.facts["pptx"]["picture_count_by_slide"] == [0, 0]
    assert observation.facts["pptx"]["granular_picture_count_by_slide"] == [0, 0]
    assert observation.facts["pptx"]["slides_with_pictures"] == 0
    assert observation.facts["pptx"]["slides_with_granular_pictures"] == 0
    assert observation.facts["pptx"]["full_slide_picture_ratio"] == 0
    assert observation.facts["pptx"]["picture_coverage_ratio_mean"] == 0


@pytest.mark.asyncio
async def test_screenshot_only_presentation_fails_editability_contract(
    tmp_path,
) -> None:
    image_path = tmp_path / "slide.png"
    pptx_path = tmp_path / "deck.pptx"
    pdf_path = tmp_path / "deck.pdf"
    _write_image(image_path, size=(1600, 900))
    _write_screenshot_pptx(pptx_path, image_path)
    _write_pdf(pdf_path, page_count=1)

    observation = await observe_creative_artifacts(
        CreativeArtifactContract(
            modality="presentation",
            aspect_ratio="16:9",
            page_count=1,
            editable_required=True,
        ),
        {"pptx": pptx_path, "pdf": pdf_path},
    )

    assert observation.hard_gates["editability"].passed is False
    assert observation.facts["pptx"]["full_slide_image_only_ratio"] == 1
    assert observation.facts["pptx"]["picture_count_by_slide"] == [1]
    assert observation.facts["pptx"]["granular_picture_count_by_slide"] == [0]
    assert observation.facts["pptx"]["distinct_picture_asset_count"] == 1
    assert observation.facts["pptx"]["full_slide_picture_ratio"] == 1
    assert observation.facts["pptx"]["slides_with_granular_pictures_ratio"] == 0
    assert observation.facts["pptx"]["slides_with_pictures_ratio"] == 1
    assert observation.facts["pptx"]["picture_coverage_ratio_by_slide"] == [1]


@pytest.mark.asyncio
async def test_missing_presentation_preview_fails_contract(tmp_path) -> None:
    pptx_path = tmp_path / "deck.pptx"
    _write_editable_pptx(pptx_path)

    observation = await observe_creative_artifacts(
        CreativeArtifactContract(
            modality="presentation",
            aspect_ratio="16:9",
            page_count=2,
            preview_required=True,
        ),
        {"pptx": pptx_path},
    )

    assert observation.hard_gates["pptx_and_preview_valid"].passed is False
    assert "pdf_missing" in observation.warnings


@pytest.mark.asyncio
async def test_image_led_presentation_picture_coverage_gate_rejects_sparse_deck(
    tmp_path,
) -> None:
    image_path = tmp_path / "asset.png"
    pptx_path = tmp_path / "deck.pptx"
    pdf_path = tmp_path / "deck.pdf"
    _write_image(image_path, size=(1600, 900))
    _write_picture_mix_pptx(pptx_path, image_path, full_slide_first=False)
    _write_pdf(pdf_path)

    observation = await observe_creative_artifacts(
        CreativeArtifactContract(
            modality="presentation",
            aspect_ratio="16:9",
            page_count=2,
            minimum_picture_coverage_ratio=0.35,
        ),
        {"pptx": pptx_path, "pdf": pdf_path},
    )

    assert observation.hard_gates["minimum_picture_coverage"].passed is False
    assert observation.facts["pptx"]["picture_coverage_ratio_mean"] < 0.35
    scenario = next(
        item
        for item in generate_evaluation_bundle(seed=20260727).manifest.public_scenarios
        if item.modality == "presentation"
    )
    evaluation = score_artifact_observation(scenario, observation)
    assert evaluation.status == "blocked"
    assert "minimum_picture_coverage" in evaluation.hard_gate_failures


@pytest.mark.asyncio
async def test_image_led_presentation_picture_coverage_gate_accepts_mean_threshold(
    tmp_path,
) -> None:
    image_path = tmp_path / "asset.png"
    pptx_path = tmp_path / "deck.pptx"
    pdf_path = tmp_path / "deck.pdf"
    _write_image(image_path, size=(1600, 900))
    _write_picture_mix_pptx(pptx_path, image_path, full_slide_first=True)
    _write_pdf(pdf_path)

    observation = await observe_creative_artifacts(
        CreativeArtifactContract(
            modality="presentation",
            aspect_ratio="16:9",
            page_count=2,
            minimum_picture_coverage_ratio=0.35,
        ),
        {"pptx": pptx_path, "pdf": pdf_path},
    )

    assert observation.hard_gates["minimum_picture_coverage"].passed is True
    assert observation.facts["pptx"]["picture_coverage_ratio_mean"] == 0.5
