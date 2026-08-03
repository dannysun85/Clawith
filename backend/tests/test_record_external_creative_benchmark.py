from __future__ import annotations

import json
from pathlib import Path

import fitz
from PIL import Image
import pytest
from pptx import Presentation
from pptx.util import Inches

from scripts.creative_provider_benchmark import BenchmarkContractError, load_case
from scripts.record_external_creative_benchmark import (
    import_external_artifact,
)


def _write_plan(path: Path, *, case_key: str, modality: str, **case: object) -> None:
    payload = {
        "benchmark_id": "external-import-test",
        "cases": {
            case_key: {
                "modality": modality,
                **case,
            }
        },
    }
    path.write_text(
        json.dumps(payload, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )


def _write_image(path: Path, *, size: tuple[int, int] = (900, 900)) -> None:
    Image.new("RGB", size, color=(24, 72, 120)).save(path, format="PNG")


def _write_presentation_pair(pptx_path: Path, pdf_path: Path) -> None:
    image_path = pptx_path.with_suffix(".png")
    _write_image(image_path, size=(640, 360))

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    for index in range(2):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(
            str(image_path),
            Inches(0.4),
            Inches(0.4),
            Inches(4.0),
            Inches(2.25),
        )
        slide.shapes.add_textbox(
            Inches(5),
            Inches(1),
            Inches(6),
            Inches(1),
        ).text = f"Slide {index + 1}"
    presentation.save(pptx_path)

    document = fitz.open()
    for _ in range(2):
        document.new_page(width=960, height=540)
    document.save(pdf_path)
    document.close()


@pytest.mark.asyncio
async def test_import_external_image_binds_provenance_and_never_generates(
    tmp_path: Path,
) -> None:
    plan_path = tmp_path / "plan.json"
    source_path = tmp_path / "doubao.png"
    output_dir = tmp_path / "out"
    _write_plan(
        plan_path,
        case_key="image_case",
        modality="image",
        prompt="same image task",
        aspect_ratio="1:1",
    )
    _write_image(source_path)
    case = load_case(plan_path, "image_case")

    receipt_path, receipt = await import_external_artifact(
        provider="doubao",
        case=case,
        artifacts={"image": source_path},
        output_dir=output_dir,
        source_reference="doubao://chat/image-case",
    )

    copied_path = Path(receipt["artifact_paths"]["image"])
    assert receipt_path.is_file()
    assert copied_path.is_file()
    assert copied_path.read_bytes() == source_path.read_bytes()
    assert receipt["provider"] == "doubao"
    assert receipt["benchmark_plan_sha256"] == case["__benchmark_plan_sha256"]
    assert receipt["benchmark_case_sha256"] == case["__benchmark_case_sha256"]
    assert receipt["cost_guardrail"]["generation_performed"] is False
    assert receipt["provider_receipt"]["acceptance_observed"] is False
    assert receipt["structural_observation"]["hard_gates"][
        "artifact_decodable"
    ]["passed"] is True
    assert copied_path.stat().st_mode & 0o777 == 0o600
    assert receipt_path.stat().st_mode & 0o777 == 0o600


@pytest.mark.asyncio
async def test_import_external_presentation_requires_both_editable_and_preview_files(
    tmp_path: Path,
) -> None:
    plan_path = tmp_path / "plan.json"
    pptx_path = tmp_path / "deck.pptx"
    _write_plan(
        plan_path,
        case_key="presentation_case",
        modality="presentation",
        goal="same deck task",
        slides=2,
        aspect_ratio="16:9",
    )
    _write_presentation_pair(pptx_path, tmp_path / "unused.pdf")
    case = load_case(plan_path, "presentation_case")

    with pytest.raises(BenchmarkContractError, match="missing_external_pdf_artifact"):
        await import_external_artifact(
            provider="doubao",
            case=case,
            artifacts={"pptx": pptx_path},
            output_dir=tmp_path / "out",
        )


@pytest.mark.asyncio
async def test_import_external_presentation_records_structural_observation(
    tmp_path: Path,
) -> None:
    plan_path = tmp_path / "plan.json"
    pptx_path = tmp_path / "deck.pptx"
    pdf_path = tmp_path / "deck.pdf"
    _write_plan(
        plan_path,
        case_key="presentation_case",
        modality="presentation",
        goal="same deck task",
        slides=2,
        aspect_ratio="16:9",
    )
    _write_presentation_pair(pptx_path, pdf_path)
    case = load_case(plan_path, "presentation_case")

    _, receipt = await import_external_artifact(
        provider="doubao",
        case=case,
        artifacts={"pptx": pptx_path, "pdf": pdf_path},
        output_dir=tmp_path / "out",
    )

    gates = receipt["structural_observation"]["hard_gates"]
    assert gates["pptx_and_preview_valid"]["passed"] is True
    assert gates["page_count_and_aspect_match"]["passed"] is True
    assert gates["editability"]["passed"] is True
    assert receipt["structural_observation"]["facts"]["pptx"]["page_count"] == 2


@pytest.mark.asyncio
async def test_import_external_video_records_structural_failure_without_success_claim(
    tmp_path: Path,
) -> None:
    plan_path = tmp_path / "plan.json"
    source_path = tmp_path / "doubao.mp4"
    _write_plan(
        plan_path,
        case_key="video_case",
        modality="video",
        prompt="same video task",
        aspect_ratio="9:16",
        duration_seconds=4,
    )
    source_path.write_bytes(b"not a decodable mp4")
    case = load_case(plan_path, "video_case")

    _, receipt = await import_external_artifact(
        provider="doubao",
        case=case,
        artifacts={"mp4": source_path},
        output_dir=tmp_path / "out",
    )

    gates = receipt["structural_observation"]["hard_gates"]
    assert gates["artifact_decodable"]["passed"] is False
    assert receipt["evidence_level"] == "external_artifact_imported"
    assert receipt["provider_receipt"]["acceptance_observed"] is False
    assert receipt["cost_guardrail"]["generation_performed"] is False


@pytest.mark.asyncio
async def test_import_external_artifact_is_provider_allowlisted_and_idempotent(
    tmp_path: Path,
) -> None:
    plan_path = tmp_path / "plan.json"
    source_path = tmp_path / "doubao.png"
    output_dir = tmp_path / "out"
    _write_plan(
        plan_path,
        case_key="image_case",
        modality="image",
        prompt="same image task",
        aspect_ratio="1:1",
    )
    _write_image(source_path)
    case = load_case(plan_path, "image_case")

    with pytest.raises(ValueError, match="unsupported external provider"):
        await import_external_artifact(
            provider="minimax",
            case=case,
            artifacts={"image": source_path},
            output_dir=output_dir,
        )

    await import_external_artifact(
        provider="doubao",
        case=case,
        artifacts={"image": source_path},
        output_dir=output_dir,
    )
    with pytest.raises(FileExistsError, match="already exists"):
        await import_external_artifact(
            provider="doubao",
            case=case,
            artifacts={"image": source_path},
            output_dir=output_dir,
        )
