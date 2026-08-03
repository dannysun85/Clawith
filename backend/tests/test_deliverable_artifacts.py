"""Authoritative Runtime artifact reconciliation for deliverable requests."""

from __future__ import annotations

import hashlib
from io import BytesIO
from types import SimpleNamespace
import uuid
import zipfile

import fitz
import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.models.agent_tool_execution import AgentToolExecution
from app.models.deliverable import DeliverableArtifactRevision, DeliverableRequest
from app.services.deliverable_artifacts import (
    DeliverableArtifactError,
    approve_deliverable_artifacts,
    deliverable_artifact_snapshot_key,
    reconcile_runtime_deliverable_artifacts,
)
from app.services.deliverable_quality_gate import (
    DeliverableQualityGateReceipt,
    attach_deliverable_quality_gate_receipt,
    blocked_quality_receipt_from_automated_evidence,
)
from app.services.storage import agent_storage_key
from app.services.storage_runtime.local import LocalStorageBackend


class _Result:
    def __init__(self, values: object) -> None:
        self.values = values

    def scalar_one_or_none(self):
        if isinstance(self.values, list):
            return self.values[0] if self.values else None
        return self.values

    def scalars(self):
        return self

    def all(self):
        return self.values if isinstance(self.values, list) else [self.values]


class _Session:
    def __init__(self, *results: object) -> None:
        self.results = list(results)
        self.added: list[DeliverableArtifactRevision] = []

    async def execute(self, _statement):
        return _Result(self.results.pop(0))

    def add(self, artifact: DeliverableArtifactRevision) -> None:
        self.added.append(artifact)


def _request() -> DeliverableRequest:
    return DeliverableRequest(
        id=uuid.uuid4(),
        tenant_id=uuid.uuid4(),
        created_by_user_id=uuid.uuid4(),
        agent_id=uuid.uuid4(),
        session_id=uuid.uuid4(),
        agent_run_id=uuid.uuid4(),
        launch_message_id=uuid.uuid4(),
        client_request_id=uuid.uuid4(),
        request_fingerprint="f" * 64,
        work_type="presentation",
        workflow_id="builtin.presentation.v1",
        workflow_version="1.0.0",
        goal="Create a launch deck",
        inputs=[],
        spec={"audience": "customers", "page_count": 8, "language": "en-US", "style": "clean"},
        tier="pro",
        approval_policy=["outline", "final"],
        output_contract=["pptx", "pdf"],
        status="running",
        current_stage="running",
        version=2,
    )


def _video_request() -> DeliverableRequest:
    request = _request()
    request.work_type = "video"
    request.workflow_id = "builtin.video.v1"
    request.goal = "Create a people-led product ad"
    request.spec = {
        "channel": "social",
        "aspect_ratio": "9:16",
        "duration": "10",
        "audience": "urban professionals",
        "language": "zh-CN",
        "audio_mode": "voiceover",
        "story": "An adult commuter uses the product",
        "cta": "Learn more",
    }
    request.approval_policy = ["storyboard", "final"]
    request.output_contract = ["mp4"]
    return request


def _poster_request() -> DeliverableRequest:
    request = _request()
    request.work_type = "poster"
    request.workflow_id = "builtin.poster.v1"
    request.goal = "Create a premium social campaign image"
    request.spec = {
        "channel": "social",
        "aspect_ratio": "16:9",
        "exact_copy": "",
        "style": "premium commercial",
    }
    request.approval_policy = ["composition", "final"]
    request.output_contract = ["png"]
    return request


def _execution(
    request: DeliverableRequest,
    *,
    tool_name: str,
    artifact_type: str,
    path: str | None = None,
) -> AgentToolExecution:
    workspace_path = path or f"workspace/deliverables/{request.id}/result.{artifact_type}"
    return AgentToolExecution(
        id=uuid.uuid4(),
        tenant_id=request.tenant_id,
        run_id=request.agent_run_id,
        tool_call_id=f"call-{artifact_type}",
        tool_name=tool_name,
        assistant_message_id=f"assistant-{artifact_type}",
        arguments_hash="a" * 64,
        sanitized_arguments={},
        effect="write",
        retry_policy="conditional",
        status="succeeded",
        result_metadata={
            "artifact_refs": [f"workspace://{request.agent_id}/{workspace_path}"],
        },
    )


def _failed_execution(
    request: DeliverableRequest,
    *,
    tool_name: str,
    artifact_type: str,
    run_id: uuid.UUID | None = None,
    error_code: str = "conversion_artifact_invalid",
) -> AgentToolExecution:
    workspace_path = f"workspace/deliverables/{request.id}/result.{artifact_type}"
    return AgentToolExecution(
        id=uuid.uuid4(),
        tenant_id=request.tenant_id,
        run_id=run_id or request.agent_run_id,
        tool_call_id=f"failed-{artifact_type}",
        tool_name=tool_name,
        assistant_message_id=f"assistant-failed-{artifact_type}",
        arguments_hash="b" * 64,
        sanitized_arguments={"target_path": workspace_path},
        effect="write",
        retry_policy="conditional",
        status="failed",
        result_summary="render quality failed",
        result_metadata={"error_code": error_code, "artifact_refs": []},
    )


def _pptx_bytes(*, page_count: int = 8, width: int = 12_192_000, height: int = 6_858_000) -> bytes:
    output = BytesIO()
    with zipfile.ZipFile(output, "w") as archive:
        archive.writestr("[Content_Types].xml", "<Types />")
        slide_ids = "".join(
            f'<p:sldId id="{256 + index}" r:id="rId{index + 1}"/>'
            for index in range(page_count)
        )
        archive.writestr(
            "ppt/presentation.xml",
            (
                '<p:presentation '
                'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
                'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
                f"<p:sldIdLst>{slide_ids}</p:sldIdLst>"
                f'<p:sldSz cx="{width}" cy="{height}"/>'
                "</p:presentation>"
            ),
        )
    return output.getvalue()


def _picture_pptx_bytes(*, sparse: bool) -> bytes:
    image_output = BytesIO()
    Image.new("RGB", (1600, 900), color=(28, 34, 46)).save(
        image_output,
        format="PNG",
    )
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    for index in range(8):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        image_output.seek(0)
        if not sparse or index == 0:
            if sparse:
                slide.shapes.add_picture(
                    image_output,
                    Inches(0.5),
                    Inches(0.5),
                    Inches(2),
                    Inches(1),
                )
            else:
                slide.shapes.add_picture(
                    image_output,
                    0,
                    0,
                    presentation.slide_width,
                    presentation.slide_height,
                )
        box = slide.shapes.add_textbox(
            Inches(1),
            Inches(1.5),
            Inches(8),
            Inches(1),
        )
        box.text = f"Slide {index + 1}"
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _pdf_bytes(*, page_count: int = 8, width: float = 960, height: float = 540) -> bytes:
    document = fitz.open()
    for _ in range(page_count):
        document.new_page(width=width, height=height)
    data = document.tobytes()
    document.close()
    return data


def _png_bytes(*, width: int = 1600, height: int = 900) -> bytes:
    output = BytesIO()
    Image.new("RGB", (width, height), color=(28, 34, 46)).save(output, format="PNG")
    return output.getvalue()


@pytest.mark.asyncio
async def test_reconcile_accepts_request_scoped_generated_png(tmp_path) -> None:
    request = _poster_request()
    storage = LocalStorageBackend(str(tmp_path))
    png_path = f"workspace/deliverables/{request.id}/final.png"
    await storage.write_bytes(agent_storage_key(request.agent_id, png_path), _png_bytes())
    db = _Session(
        [_execution(request, tool_name="generate_image_minimax", artifact_type="png", path=png_path)],
        [],
    )

    result = await reconcile_runtime_deliverable_artifacts(
        db,  # type: ignore[arg-type]
        request=request,
        run_id=request.agent_run_id,
        storage=storage,
    )

    assert result.complete is True
    assert len(result.artifacts) == 1
    artifact = result.artifacts[0]
    assert artifact.artifact_type == "png"
    assert artifact.mime_type == "image/png"
    assert artifact.evaluation["verified"] is True
    assert artifact.evaluation["verification_level"] == "contract"


@pytest.mark.asyncio
async def test_reconcile_accepts_request_scoped_browser_safe_video_with_voiceover(
    tmp_path,
    monkeypatch,
) -> None:
    request = _video_request()
    storage = LocalStorageBackend(str(tmp_path))
    mp4_path = f"workspace/deliverables/{request.id}/final.mp4"
    raw = b"\x00\x00\x00\x18ftypmp42" + b"video-payload"
    await storage.write_bytes(agent_storage_key(request.agent_id, mp4_path), raw)
    info = SimpleNamespace(
        width=720,
        height=1280,
        duration_seconds=10.0,
        codec_name="h264",
        pixel_format="yuv420p",
        audio_codec_name="aac",
        fast_start=True,
    )

    async def validate_video(*_args, **_kwargs):
        return info

    monkeypatch.setattr(
        "app.services.media_assets.validate_generated_video",
        validate_video,
    )
    monkeypatch.setattr(
        "app.services.media_assets.validate_video_delivery_contract",
        lambda value, **_kwargs: value,
    )
    db = _Session(
        [_execution(request, tool_name="compose_video_audio", artifact_type="mp4", path=mp4_path)],
        [],
    )

    result = await reconcile_runtime_deliverable_artifacts(
        db,  # type: ignore[arg-type]
        request=request,
        run_id=request.agent_run_id,
        storage=storage,
    )

    assert result.complete is True
    assert len(result.artifacts) == 1
    artifact = result.artifacts[0]
    assert artifact.artifact_type == "mp4"
    assert artifact.mime_type == "video/mp4"
    assert artifact.evaluation["facts"]["audio_codec"] == "aac"
    assert artifact.evaluation["facts"]["audio_mode"] == "voiceover"
    assert "browser_codec" in artifact.evaluation["checks"]


@pytest.mark.asyncio
async def test_reconcile_persists_only_exact_request_scoped_structural_outputs(tmp_path) -> None:
    request = _request()
    storage = LocalStorageBackend(str(tmp_path))
    pptx_path = f"workspace/deliverables/{request.id}/result.pptx"
    pdf_path = f"workspace/deliverables/{request.id}/result.pdf"
    await storage.write_bytes(agent_storage_key(request.agent_id, pptx_path), _pptx_bytes())
    await storage.write_bytes(
        agent_storage_key(request.agent_id, pdf_path),
        _pdf_bytes(),
    )
    executions = [
        _execution(request, tool_name="convert_html_to_pptx", artifact_type="pptx"),
        _execution(request, tool_name="convert_html_to_pdf", artifact_type="pdf"),
        _execution(
            request,
            tool_name="convert_html_to_pdf",
            artifact_type="pdf",
            path="workspace/deliverables/another-request/foreign.pdf",
        ),
    ]
    db = _Session(executions, [])

    result = await reconcile_runtime_deliverable_artifacts(
        db,  # type: ignore[arg-type]
        request=request,
        run_id=request.agent_run_id,
        storage=storage,
    )

    assert result.complete is True
    assert {artifact.artifact_key for artifact in result.artifacts} == {"pptx", "pdf"}
    assert {artifact.workspace_path for artifact in result.artifacts} == {pptx_path, pdf_path}
    assert len(db.added) == 2
    assert all(artifact.evaluation["verified"] is True for artifact in db.added)
    assert all(artifact.evaluation["verification_level"] == "contract" for artifact in db.added)
    assert all(artifact.evaluation["facts"]["page_count"] == 8 for artifact in db.added)
    for artifact in db.added:
        snapshot = await storage.read_bytes(deliverable_artifact_snapshot_key(artifact))
        assert artifact.content_hash == hashlib.sha256(snapshot).hexdigest()


@pytest.mark.asyncio
async def test_reconcile_rejects_sparse_image_led_presentation(tmp_path) -> None:
    request = _request()
    request.goal = "制作一份图文并茂的新品发布方案"
    storage = LocalStorageBackend(str(tmp_path))
    pptx_path = f"workspace/deliverables/{request.id}/result.pptx"
    pdf_path = f"workspace/deliverables/{request.id}/result.pdf"
    await storage.write_bytes(
        agent_storage_key(request.agent_id, pptx_path),
        _picture_pptx_bytes(sparse=True),
    )
    await storage.write_bytes(
        agent_storage_key(request.agent_id, pdf_path),
        _pdf_bytes(),
    )

    result = await reconcile_runtime_deliverable_artifacts(
        _Session(
            [
                _execution(request, tool_name="convert_html_to_pptx", artifact_type="pptx"),
                _execution(request, tool_name="convert_html_to_pdf", artifact_type="pdf"),
            ],
            [],
        ),
        request=request,
        run_id=request.agent_run_id,
        storage=storage,
    )

    assert result.complete is False
    assert result.invalid_types == ("pptx", "pdf")
    assert result.failure_codes == (
        ("pptx", "presentation_picture_coverage_below_minimum"),
        ("pdf", "presentation_picture_coverage_below_minimum"),
    )
    assert result.artifacts == ()


@pytest.mark.asyncio
async def test_reconcile_accepts_image_led_presentation_with_meaningful_coverage(
    tmp_path,
) -> None:
    request = _request()
    request.goal = "制作一份图文并茂的新品发布方案"
    storage = LocalStorageBackend(str(tmp_path))
    pptx_path = f"workspace/deliverables/{request.id}/result.pptx"
    pdf_path = f"workspace/deliverables/{request.id}/result.pdf"
    await storage.write_bytes(
        agent_storage_key(request.agent_id, pptx_path),
        _picture_pptx_bytes(sparse=False),
    )
    await storage.write_bytes(
        agent_storage_key(request.agent_id, pdf_path),
        _pdf_bytes(),
    )

    result = await reconcile_runtime_deliverable_artifacts(
        _Session(
            [
                _execution(request, tool_name="convert_html_to_pptx", artifact_type="pptx"),
                _execution(request, tool_name="convert_html_to_pdf", artifact_type="pdf"),
            ],
            [],
        ),
        request=request,
        run_id=request.agent_run_id,
        storage=storage,
    )

    assert result.complete is True
    assert result.invalid_types == ()
    pptx_artifact = next(item for item in result.artifacts if item.artifact_type == "pptx")
    assert pptx_artifact.evaluation["facts"]["picture_coverage_ratio_mean"] == 1
    assert pptx_artifact.evaluation["facts"]["picture_coverage_gate"] == 1


@pytest.mark.asyncio
async def test_reconcile_fails_closed_when_required_pdf_is_invalid(tmp_path) -> None:
    request = _request()
    storage = LocalStorageBackend(str(tmp_path))
    pptx_path = f"workspace/deliverables/{request.id}/result.pptx"
    pdf_path = f"workspace/deliverables/{request.id}/result.pdf"
    await storage.write_bytes(agent_storage_key(request.agent_id, pptx_path), _pptx_bytes())
    await storage.write_bytes(agent_storage_key(request.agent_id, pdf_path), b"not a pdf")
    db = _Session(
        [
            _execution(request, tool_name="convert_html_to_pptx", artifact_type="pptx"),
            _execution(request, tool_name="convert_html_to_pdf", artifact_type="pdf"),
        ],
        [],
    )

    result = await reconcile_runtime_deliverable_artifacts(
        db,  # type: ignore[arg-type]
        request=request,
        run_id=request.agent_run_id,
        storage=storage,
    )

    assert result.complete is False
    assert result.invalid_types == ("pdf",)
    assert [artifact.artifact_key for artifact in db.added] == ["pptx"]


@pytest.mark.asyncio
async def test_reconcile_reports_current_run_failed_conversion_attempt(tmp_path) -> None:
    request = _request()
    storage = LocalStorageBackend(str(tmp_path))
    failed = _failed_execution(
        request,
        tool_name="convert_html_to_pdf",
        artifact_type="pdf",
        error_code="presentation_visual_quality_failed",
    )
    db = _Session([failed], [])

    result = await reconcile_runtime_deliverable_artifacts(
        db,  # type: ignore[arg-type]
        request=request,
        run_id=request.agent_run_id,
        storage=storage,
    )

    assert result.complete is False
    assert result.attempted_types == ("pdf",)
    assert result.failed_types == ("pdf",)
    assert result.invalid_types == ("pdf",)
    assert result.failure_codes == (
        ("pdf", "presentation_visual_quality_failed"),
    )


@pytest.mark.asyncio
async def test_reconcile_repairs_one_type_and_reuses_other_verified_candidate(
    tmp_path,
) -> None:
    request = _request()
    storage = LocalStorageBackend(str(tmp_path))
    pptx_path = f"workspace/deliverables/{request.id}/result.pptx"
    pdf_path = f"workspace/deliverables/{request.id}/result.pdf"
    await storage.write_bytes(agent_storage_key(request.agent_id, pptx_path), _pptx_bytes())
    await storage.write_bytes(agent_storage_key(request.agent_id, pdf_path), _pdf_bytes())
    initial = await reconcile_runtime_deliverable_artifacts(
        _Session(
            [
                _execution(request, tool_name="convert_html_to_pptx", artifact_type="pptx"),
                _execution(request, tool_name="convert_html_to_pdf", artifact_type="pdf"),
            ],
            [],
        ),  # type: ignore[arg-type]
        request=request,
        run_id=request.agent_run_id,
        storage=storage,
    )
    followup_run_id = uuid.uuid4()
    await storage.write_bytes(
        agent_storage_key(request.agent_id, pdf_path),
        _pdf_bytes(width=1280, height=720),
    )
    repaired_pdf = _execution(
        request,
        tool_name="convert_html_to_pdf",
        artifact_type="pdf",
    )
    repaired_pdf.run_id = followup_run_id
    repaired = await reconcile_runtime_deliverable_artifacts(
        _Session([repaired_pdf], list(initial.artifacts)),  # type: ignore[arg-type]
        request=request,
        run_id=followup_run_id,
        storage=storage,
    )

    assert repaired.complete is True
    assert repaired.attempted_types == ("pdf",)
    assert repaired.created_types == ("pdf",)
    assert {artifact.artifact_type for artifact in repaired.artifacts} == {
        "pptx",
        "pdf",
    }
    latest_pdf = next(
        artifact for artifact in repaired.artifacts if artifact.artifact_type == "pdf"
    )
    assert latest_pdf.revision_number == 2


@pytest.mark.asyncio
async def test_reconcile_rejects_pdf_with_wrong_page_count(tmp_path) -> None:
    request = _request()
    storage = LocalStorageBackend(str(tmp_path))
    pptx_path = f"workspace/deliverables/{request.id}/result.pptx"
    pdf_path = f"workspace/deliverables/{request.id}/result.pdf"
    await storage.write_bytes(agent_storage_key(request.agent_id, pptx_path), _pptx_bytes())
    await storage.write_bytes(agent_storage_key(request.agent_id, pdf_path), _pdf_bytes(page_count=5))
    db = _Session(
        [
            _execution(request, tool_name="convert_html_to_pptx", artifact_type="pptx"),
            _execution(request, tool_name="convert_html_to_pdf", artifact_type="pdf"),
        ],
        [],
    )

    result = await reconcile_runtime_deliverable_artifacts(
        db,  # type: ignore[arg-type]
        request=request,
        run_id=request.agent_run_id,
        storage=storage,
    )

    assert result.complete is False
    assert set(result.invalid_types) == {"pptx", "pdf"}
    assert db.added == []


@pytest.mark.asyncio
async def test_reconcile_rejects_non_widescreen_pdf(tmp_path) -> None:
    request = _request()
    storage = LocalStorageBackend(str(tmp_path))
    pptx_path = f"workspace/deliverables/{request.id}/result.pptx"
    pdf_path = f"workspace/deliverables/{request.id}/result.pdf"
    await storage.write_bytes(agent_storage_key(request.agent_id, pptx_path), _pptx_bytes())
    await storage.write_bytes(
        agent_storage_key(request.agent_id, pdf_path),
        _pdf_bytes(width=1008, height=576),
    )
    db = _Session(
        [
            _execution(request, tool_name="convert_html_to_pptx", artifact_type="pptx"),
            _execution(request, tool_name="convert_html_to_pdf", artifact_type="pdf"),
        ],
        [],
    )

    result = await reconcile_runtime_deliverable_artifacts(
        db,  # type: ignore[arg-type]
        request=request,
        run_id=request.agent_run_id,
        storage=storage,
    )

    assert result.complete is False
    assert result.invalid_types == ("pdf",)
    assert [artifact.artifact_key for artifact in db.added] == ["pptx"]


@pytest.mark.asyncio
async def test_approval_rechecks_content_hash_before_accepting_mutable_workspace_path(tmp_path) -> None:
    request = _request()
    storage = LocalStorageBackend(str(tmp_path))
    pptx_path = f"workspace/deliverables/{request.id}/result.pptx"
    pdf_path = f"workspace/deliverables/{request.id}/result.pdf"
    pdf_data = _pdf_bytes()
    await storage.write_bytes(agent_storage_key(request.agent_id, pptx_path), _pptx_bytes())
    await storage.write_bytes(agent_storage_key(request.agent_id, pdf_path), pdf_data)
    reconcile_db = _Session(
        [
            _execution(request, tool_name="convert_html_to_pptx", artifact_type="pptx"),
            _execution(request, tool_name="convert_html_to_pdf", artifact_type="pdf"),
        ],
        [],
    )
    reconciled = await reconcile_runtime_deliverable_artifacts(
        reconcile_db,  # type: ignore[arg-type]
        request=request,
        run_id=request.agent_run_id,
        storage=storage,
    )

    approved = await approve_deliverable_artifacts(
        _Session(list(reconciled.artifacts)),  # type: ignore[arg-type]
        request=request,
        storage=storage,
    )
    assert {artifact.artifact_key for artifact in approved} == {"pptx", "pdf"}

    with pytest.raises(DeliverableArtifactError) as required_quality_error:
        await approve_deliverable_artifacts(
            _Session(list(reconciled.artifacts)),  # type: ignore[arg-type]
            request=request,
            storage=storage,
            require_creative_quality_gate=True,
        )
    assert (
        required_quality_error.value.code
        == "deliverable_creative_quality_review_required"
    )

    artifact_hashes = {
        artifact.artifact_key: artifact.content_hash
        for artifact in reconciled.artifacts
    }
    attach_deliverable_quality_gate_receipt(
        reconciled.artifacts,
        blocked_quality_receipt_from_automated_evidence(
            receipt_ref="ocr:blocking-finding",
            artifact_hashes=artifact_hashes,
            evidence_kind="ocr",
            hard_gate_failures=("no_unrequested_watermark",),
        ),
    )
    with pytest.raises(DeliverableArtifactError) as blocked_quality_error:
        await approve_deliverable_artifacts(
            _Session(list(reconciled.artifacts)),  # type: ignore[arg-type]
            request=request,
            storage=storage,
            require_creative_quality_gate=False,
        )
    assert blocked_quality_error.value.code == "deliverable_creative_quality_blocked"

    attach_deliverable_quality_gate_receipt(
        reconciled.artifacts,
        DeliverableQualityGateReceipt(
            receipt_ref="panel:commercial-pass",
            source="blind_review_panel",
            status="passed",
            artifact_hashes=artifact_hashes,
            reviewer_count=3,
            required_evidence_kinds=("document_semantic", "human_visual"),
            complete_evidence_kinds=("document_semantic", "human_visual"),
            commercially_usable=True,
        ),
    )
    previous_approved = tuple(
        DeliverableArtifactRevision(
            id=uuid.uuid4(),
            tenant_id=request.tenant_id,
            request_id=request.id,
            artifact_key=artifact.artifact_key,
            artifact_type=artifact.artifact_type,
            workspace_path=f"workspace/deliverables/{request.id}/previous.{artifact.artifact_type}",
            mime_type=artifact.mime_type,
            content_hash=("c" if artifact.artifact_type == "pptx" else "d") * 64,
            size_bytes=artifact.size_bytes,
            revision_number=1,
            status="approved",
            evaluation={"verified": True},
        )
        for artifact in reconciled.artifacts
    )
    for artifact in reconciled.artifacts:
        artifact.revision_number = 2
    approved_with_quality = await approve_deliverable_artifacts(
        _Session([*reconciled.artifacts, *previous_approved]),  # type: ignore[arg-type]
        request=request,
        storage=storage,
        require_creative_quality_gate=True,
    )
    assert {artifact.artifact_key for artifact in approved_with_quality} == {
        "pptx",
        "pdf",
    }
    assert {artifact.status for artifact in previous_approved} == {"superseded"}

    request.spec = {**request.spec, "page_count": 7}
    with pytest.raises(DeliverableArtifactError) as contract_error:
        await approve_deliverable_artifacts(
            _Session(list(reconciled.artifacts)),  # type: ignore[arg-type]
            request=request,
            storage=storage,
        )
    assert contract_error.value.code == "deliverable_artifact_contract_invalid"
    request.spec = {**request.spec, "page_count": 8}

    for artifact in reconciled.artifacts:
        assert await storage.exists(deliverable_artifact_snapshot_key(artifact))

    await storage.write_bytes(agent_storage_key(request.agent_id, pdf_path), b"%PDF-1.7\nchanged\n%%EOF")
    with pytest.raises(DeliverableArtifactError) as error:
        await approve_deliverable_artifacts(
            _Session(list(reconciled.artifacts)),  # type: ignore[arg-type]
            request=request,
            storage=storage,
        )
    assert error.value.code == "deliverable_artifact_changed"

    await storage.write_bytes(agent_storage_key(request.agent_id, pdf_path), pdf_data)
    pdf_artifact = next(item for item in reconciled.artifacts if item.artifact_type == "pdf")
    await storage.write_bytes(deliverable_artifact_snapshot_key(pdf_artifact), b"tampered snapshot")
    with pytest.raises(DeliverableArtifactError) as snapshot_error:
        await approve_deliverable_artifacts(
            _Session(list(reconciled.artifacts)),  # type: ignore[arg-type]
            request=request,
            storage=storage,
        )
    assert snapshot_error.value.code == "deliverable_artifact_snapshot_changed"
