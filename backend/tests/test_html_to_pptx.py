"""Regression coverage for high-fidelity HTML-to-PPTX conversion."""

import pytest
from PIL import Image
from pptx import Presentation

from app.services.document_conversion import pptx_renderer


def test_slide_size_preserves_portrait_source_aspect_ratio() -> None:
    width, height = pptx_renderer._slide_size_inches(1080, 1920)

    assert width == pytest.approx(7.5)
    assert height == pytest.approx(13.333333)


@pytest.mark.asyncio
async def test_default_visual_conversion_preserves_browser_slide_and_cleans_temp_file(monkeypatch, tmp_path):
    src_file = tmp_path / "source.html"
    tgt_file = tmp_path / "result.pptx"
    screenshot = tmp_path / "browser-slide.png"
    src_file.write_text(
        "<html><body><section><h1>Title</h1><table><tr><td>Cell</td></tr></table></section></body></html>",
        encoding="utf-8",
    )
    Image.new("RGB", (1280, 720), color=(24, 36, 52)).save(screenshot)

    async def fake_collect_browser_layout(_src, width, height, render_mode, render_scale):
        assert (width, height, render_mode, render_scale) == (1280, 720, "visual", 2.0)
        return {
            "slides": [{"width": 1280, "height": 720, "backgroundColor": "#ffffff"}],
            "screenshots": [str(screenshot)],
        }

    monkeypatch.setattr(pptx_renderer, "collect_browser_layout", fake_collect_browser_layout)

    result = await pptx_renderer.render_html_to_pptx(
        src_file,
        tgt_file,
        "workspace/result.pptx",
        tmp_path,
        {},
    )

    assert result.startswith("✅ Successfully converted HTML to high-fidelity PPTX screenshots")
    assert tgt_file.exists()
    presentation = Presentation(tgt_file)
    assert len(presentation.slides) == 1
    assert len(presentation.slides[0].shapes) == 1
    assert presentation.slides[0].shapes[0].shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
    assert not screenshot.exists()


@pytest.mark.asyncio
async def test_portrait_visual_conversion_uses_portrait_pptx_canvas(
    monkeypatch,
    tmp_path,
):
    src_file = tmp_path / "poster.html"
    tgt_file = tmp_path / "poster.pptx"
    screenshot = tmp_path / "poster.png"
    src_file.write_text(
        '<html><body><section class="slide"></section></body></html>',
        encoding="utf-8",
    )
    Image.new("RGB", (1080, 1920), color=(45, 35, 95)).save(screenshot)

    async def fake_collect_browser_layout(
        _src,
        width,
        height,
        render_mode,
        render_scale,
    ):
        assert (width, height, render_mode, render_scale) == (
            1080,
            1920,
            "visual",
            2.0,
        )
        return {
            "slides": [{"width": 1080, "height": 1920}],
            "screenshots": [str(screenshot)],
        }

    monkeypatch.setattr(
        pptx_renderer,
        "collect_browser_layout",
        fake_collect_browser_layout,
    )

    await pptx_renderer.render_html_to_pptx(
        src_file,
        tgt_file,
        "workspace/poster.pptx",
        tmp_path,
        {"design_width": 1080, "design_height": 1920},
    )

    presentation = Presentation(tgt_file)
    assert presentation.slide_width / 914400 == pytest.approx(7.5, abs=1e-5)
    assert presentation.slide_height / 914400 == pytest.approx(13.333333, abs=1e-5)
    picture = presentation.slides[0].shapes[0]
    assert picture.width == presentation.slide_width
    assert abs(picture.height - presentation.slide_height) <= 1


@pytest.mark.asyncio
async def test_hybrid_editable_uses_visual_layer_and_editable_text(
    monkeypatch,
    tmp_path,
):
    src_file = tmp_path / "source.html"
    tgt_file = tmp_path / "result.pptx"
    content_screenshot = tmp_path / "content-without-text.png"
    src_file.write_text(
        '<html><body><section class="slide"><h1>可编辑标题</h1></section></body></html>',
        encoding="utf-8",
    )
    Image.new("RGB", (1280, 720), color=(24, 36, 52)).save(content_screenshot)

    async def fake_collect_browser_layout(
        _src,
        width,
        height,
        render_mode,
        render_scale,
    ):
        assert (width, height, render_mode, render_scale) == (
            1280,
            720,
            "hybrid_editable",
            2.0,
        )
        return {
            "slides": [
                {
                    "width": 1280,
                    "height": 720,
                    "backgroundColor": "#ffffff",
                    "items": [
                        {
                            "kind": "text",
                            "tag": "h1",
                            "text": "可编辑标题",
                            "x": 80,
                            "y": 60,
                            "w": 500,
                            "h": 80,
                            "style": {
                                "color": "rgb(255, 255, 255)",
                                "fontSize": "48px",
                            },
                        },
                        {
                            "kind": "shape",
                            "x": 60,
                            "y": 40,
                            "w": 600,
                            "h": 120,
                            "style": {"backgroundColor": "rgb(0, 0, 0)"},
                        },
                    ],
                }
            ],
            "contentScreenshots": [str(content_screenshot)],
        }

    monkeypatch.setattr(
        pptx_renderer,
        "collect_browser_layout",
        fake_collect_browser_layout,
    )

    result = await pptx_renderer.render_html_to_pptx(
        src_file,
        tgt_file,
        "workspace/result.pptx",
        tmp_path,
        {"render_mode": "hybrid_editable"},
    )

    assert result.startswith("✅ Successfully converted HTML to editable PPTX")
    presentation = Presentation(tgt_file)
    assert len(presentation.slides) == 1
    assert len(presentation.slides[0].shapes) == 2
    assert presentation.slides[0].shapes[0].shape_type != 13
    text_shapes = [
        shape
        for shape in presentation.slides[0].shapes
        if shape.has_text_frame and shape.text.strip()
    ]
    assert len(text_shapes) == 1
    assert text_shapes[0].text == "可编辑标题"
    run_xml = text_shapes[0].text_frame.paragraphs[0].runs[0]._r.xml
    assert '<a:ea typeface="Noto Sans CJK SC"' in run_xml
    assert '<a:cs typeface="Noto Sans CJK SC"' in run_xml
    assert not content_screenshot.exists()


@pytest.mark.asyncio
async def test_hybrid_editable_materializes_measured_image_instead_of_full_slide_screenshot(
    monkeypatch,
    tmp_path,
):
    src_file = tmp_path / "source.html"
    tgt_file = tmp_path / "result.pptx"
    image = tmp_path / "product.png"
    content_screenshot = tmp_path / "content-without-text.png"
    src_file.write_text(
        '<html><body><section class="slide"><img src="product.png"></section></body></html>',
        encoding="utf-8",
    )
    Image.new("RGB", (640, 480), color=(180, 80, 40)).save(image)
    Image.new("RGB", (1280, 720), color=(24, 36, 52)).save(content_screenshot)

    async def fake_collect_browser_layout(
        _src,
        width,
        height,
        render_mode,
        render_scale,
    ):
        assert (width, height, render_mode, render_scale) == (
            1280,
            720,
            "hybrid_editable",
            2.0,
        )
        return {
            "slides": [
                {
                    "width": 1280,
                    "height": 720,
                    "backgroundColor": "#ffffff",
                    "items": [
                        {
                            "kind": "image",
                            "src": "product.png",
                            "x": 120,
                            "y": 100,
                            "w": 420,
                            "h": 315,
                            "style": {},
                        },
                    ],
                }
            ],
            "contentScreenshots": [str(content_screenshot)],
        }

    monkeypatch.setattr(
        pptx_renderer,
        "collect_browser_layout",
        fake_collect_browser_layout,
    )

    result = await pptx_renderer.render_html_to_pptx(
        src_file,
        tgt_file,
        "workspace/result.pptx",
        tmp_path,
        {"render_mode": "hybrid_editable"},
    )

    assert result.startswith("✅ Successfully converted HTML to editable PPTX")
    presentation = Presentation(tgt_file)
    pictures = [
        shape
        for shape in presentation.slides[0].shapes
        if shape.shape_type == 13
    ]
    assert len(pictures) == 1
    assert pictures[0].width < presentation.slide_width * 0.6
    assert pictures[0].height < presentation.slide_height * 0.6
    assert not content_screenshot.exists()
