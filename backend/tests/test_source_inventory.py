"""FR-P2/P4/P5/P8: source inventory, semantic gate, and v2 deck quality gates."""

from __future__ import annotations

import json
from pathlib import Path
import uuid

import pytest

from app.models.deliverable import DeliverableRequest
from app.services.deliverable_artifacts import (
    _presentation_contract_facts,
    _VerifiedArtifact,
)
from app.services.document_conversion.presentation_contract import (
    PresentationVisualQualityError,
    validate_browser_slide_text_bounds,
    validate_browser_slide_visual_quality,
    validate_presentation_html_contract,
)
from app.services.source_inventory import (
    SourceInventoryEntry,
    compile_source_inventory,
    detect_fact_assertions,
    inventory_sha256,
    reconcile_slide_semantics,
    resolve_source_ref,
)


def _request(*, spec: dict | None = None, inputs: list | None = None, **overrides) -> DeliverableRequest:
    request = DeliverableRequest(
        id=uuid.uuid4(),
        tenant_id=uuid.uuid4(),
        created_by_user_id=uuid.uuid4(),
        agent_id=uuid.uuid4(),
        session_id=uuid.uuid4(),
        client_request_id=uuid.uuid4(),
        request_fingerprint="a" * 64,
        work_type="presentation",
        workflow_id="builtin.presentation.v2",
        workflow_version="2.0.0",
        goal="为极光保温杯制作客户提案演示",
        inputs=inputs or [],
        spec=spec
        or {
            "audience": "潜在客户采购负责人",
            "scenario": "client_proposal",
            "page_count": 8,
            "language": "zh-CN",
            "style": "professional",
            "key_points": "保温时长 24 小时\n内胆采用 316 不锈钢",
            "editability_contract": "editable",
            "fallback_policy": "primary_only",
        },
        tier="pro",
        approval_policy=["outline", "final"],
        output_contract=["pptx"],
        status="ready",
        current_stage="brief_confirmed",
        version=1,
        contract_revision=1,
    )
    for key, value in overrides.items():
        setattr(request, key, value)
    return request


class _FakeStorage:
    def __init__(self, payloads: dict[str, bytes] | None = None) -> None:
        self.payloads = dict(payloads or {})

    async def read_bytes(self, key: str) -> bytes:
        for suffix, payload in self.payloads.items():
            if key.endswith(suffix):
                return payload
        raise FileNotFoundError(key)

    async def read_text(self, key: str, **_kwargs) -> str:
        return (await self.read_bytes(key)).decode("utf-8")


def _pdf_with_text(text: str) -> bytes:
    import fitz

    document = fitz.open()
    page = document.new_page(width=960, height=540)
    page.insert_text((72, 96), text)
    data = document.tobytes()
    document.close()
    return data


# ─── inventory registration ─────────────────────────────────────


@pytest.mark.asyncio
async def test_inventory_registers_upload_url_and_brief_with_hash() -> None:
    pdf_bytes = _pdf_with_text("保温时长 24 小时，实测数据")
    request = _request(
        inputs=[{"path": "workspace/uploads/thermos-spec.pdf"}],
        spec={
            "audience": "潜在客户采购负责人",
            "scenario": "client_proposal",
            "page_count": 8,
            "language": "zh-CN",
            "key_points": "保温时长 24 小时",
            "source_urls": [
                {"url": "https://example.com/report", "facts": ["行业规模 35 亿元"]}
            ],
        },
    )
    entries = await compile_source_inventory(
        request,
        storage=_FakeStorage({"thermos-spec.pdf": pdf_bytes}),
    )
    kinds = [entry.kind for entry in entries]
    assert kinds == ["upload", "url", "brief"]
    upload = entries[0]
    assert upload.path == "workspace/uploads/thermos-spec.pdf"
    assert len(upload.sha256) == 64
    assert any("24" in fact for fact in upload.extracted_facts)
    url_entry = entries[1]
    assert url_entry.url == "https://example.com/report"
    assert url_entry.extracted_facts == ("行业规模 35 亿元",)
    brief_entry = entries[2]
    assert any("保温时长 24 小时" in fact for fact in brief_entry.extracted_facts)
    # The registry is hash-bound and reproducible.
    assert inventory_sha256(entries) == inventory_sha256(entries)


@pytest.mark.asyncio
async def test_inventory_unreadable_upload_registers_without_facts() -> None:
    request = _request(inputs=[{"path": "workspace/uploads/missing.pdf"}])
    entries = await compile_source_inventory(request, storage=_FakeStorage())
    upload = entries[0]
    assert upload.kind == "upload"
    assert upload.sha256 == ""
    assert upload.extracted_facts == ()
    # A fact-less entry can never source an assertion.
    assert resolve_source_ref("src-01", entries) is upload


# ─── fact assertion detection ───────────────────────────────────


def test_fact_assertion_detection_quantified_ranking_assumption() -> None:
    hits = detect_fact_assertions("保温时长可达 24 小时。市场份额排名第一。")
    assert len(hits) == 2
    assert {hit.kind for hit in hits} == {"quantified", "ranking"}
    assert all(not hit.assumption for hit in hits)

    assumed = detect_fact_assertions("假设：转化率有望提升 30%。")
    assert len(assumed) == 1
    assert assumed[0].assumption is True

    assert detect_fact_assertions("设计简洁，适合日常通勤使用。") == ()


# ─── slide semantic reconciliation ──────────────────────────────


def _entries() -> tuple[SourceInventoryEntry, ...]:
    return (
        SourceInventoryEntry(
            source_id="src-01",
            kind="brief",
            sha256="b" * 64,
            extracted_facts=("保温时长 24 小时", "内胆采用 316 不锈钢"),
            registered_by="user",
        ),
    )


def _slide(**overrides) -> dict:
    slide = {
        "slide_id": "slide-01",
        "headline": "产品实力",
        "body_points": ["保温时长 24 小时"],
        "source_refs": ["src-01"],
        "visual_kind": "editable_chart",
        "slide_type": "content",
        "layout": "metrics",
        "visual_asset": "柱状图",
        "asset_ref": "",
    }
    slide.update(overrides)
    return slide


def test_reconcile_accepts_sourced_facts() -> None:
    reconciliation = reconcile_slide_semantics([_slide()], _entries())
    assert reconciliation.passed
    assert reconciliation.assertion_count == 1
    assert len(reconciliation.inventory_sha256) == 64


def test_reconcile_blocks_unsourced_fact_assertion() -> None:
    slide = _slide(body_points=["市场占有率 48%，行业第一"])
    reconciliation = reconcile_slide_semantics([slide], _entries())
    assert not reconciliation.passed
    codes = {finding.code for finding in reconciliation.findings}
    assert "unsourced_fact_assertion" in codes


def test_reconcile_allows_explicit_assumption_label() -> None:
    slide = _slide(body_points=["假设：转化率有望提升 30%，待试点验证"])
    reconciliation = reconcile_slide_semantics([slide], _entries())
    assert reconciliation.passed
    assert reconciliation.assumption_count == 1


def test_reconcile_blocks_unresolved_source_ref() -> None:
    slide = _slide(source_refs=["src-99"])
    reconciliation = reconcile_slide_semantics([slide], _entries())
    assert not reconciliation.passed
    assert reconciliation.findings[0].code == "unresolved_source_ref"


def test_reconcile_forbids_fact_assertions_on_image_slides() -> None:
    slide = _slide(
        visual_kind="generated_image",
        asset_ref="assets/hero.png",
        body_points=["保温时长 24 小时"],
    )
    reconciliation = reconcile_slide_semantics([slide], _entries())
    assert not reconciliation.passed
    assert reconciliation.findings[0].code == "image_slide_fact_assertion"


def test_reconcile_requires_editable_visual_on_data_slides() -> None:
    slide = _slide(data_slide=True, visual_kind="editable_table")
    assert reconcile_slide_semantics([slide], _entries()).passed
    rasterized = _slide(
        data_slide=True,
        visual_kind="generated_image",
        asset_ref="assets/chart.png",
        body_points=["定性描述"],
    )
    reconciliation = reconcile_slide_semantics([rasterized], _entries())
    assert not reconciliation.passed
    assert reconciliation.findings[0].code == "data_slide_not_editable"


# ─── conversion-time semantic hard gate ─────────────────────────


def _semantic_fixture(tmp_path: Path, *, body_point: str, source_refs: list) -> tuple[Path, Path, Path]:
    source = tmp_path / "slides.html"
    source.write_text(
        """
        <html><head><style>.slide{width:1280px;height:720px}</style></head><body>
          <section class="slide" data-slide="slide-01" data-layout="cover">
            <h1 data-slide-title>产品实力</h1>
            <div data-visual><span>图表</span><span>表格</span></div>
          </section>
        </body></html>
        """,
        encoding="utf-8",
    )
    outline = tmp_path / "outline.json"
    outline.write_text(
        json.dumps(
            {
                "deck_title": "产品实力",
                "audience": "采购负责人",
                "core_message": "用可验证数据说话",
                "slides": [
                    {
                        "slide_id": "slide-01",
                        "purpose": "建立信任",
                        "headline": "产品实力",
                        "evidence": ["保温时长 24 小时"],
                        "visual_intent": "图表",
                    }
                ],
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    slide_spec = tmp_path / "slide_spec.json"
    slide_spec.write_text(
        json.dumps(
            {
                "visual_plan_version": "adaptive-v1",
                "visual_policy": {
                    "minimum_distinct_layouts": 1,
                    "minimum_distinct_images": 0,
                    "minimum_image_slides": 0,
                    "maximum_uses_per_image": 0,
                    "minimum_editable_compositions": 1,
                },
                "slides": [
                    {
                        "slide_id": "slide-01",
                        "headline": "产品实力",
                        "slide_type": "content",
                        "layout": "cover",
                        "body_points": [body_point],
                        "visual_kind": "editable_chart",
                        "visual_asset": "图表",
                        "asset_ref": "",
                        "source_refs": source_refs,
                    }
                ],
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    return source, outline, slide_spec


def test_semantic_gate_blocks_unsourced_fact_at_conversion(tmp_path: Path) -> None:
    source, outline, slide_spec = _semantic_fixture(
        tmp_path,
        body_point="市场占有率 48%",
        source_refs=["src-01"],
    )
    entries = [entry.model_dump(mode="json") for entry in _entries()]
    with pytest.raises(ValueError, match="semantic gate"):
        validate_presentation_html_contract(
            source,
            expected_page_count=1,
            outline_file=outline,
            slide_spec_file=slide_spec,
            require_adaptive_visual_plan=True,
            source_inventory_entries=entries,
            semantic_gate=True,
        )


def test_semantic_gate_accepts_sourced_fact_and_assumption(tmp_path: Path) -> None:
    entries = [entry.model_dump(mode="json") for entry in _entries()]
    sourced, outline, slide_spec = _semantic_fixture(
        tmp_path,
        body_point="保温时长 24 小时",
        source_refs=["src-01"],
    )
    validate_presentation_html_contract(
        sourced,
        expected_page_count=1,
        outline_file=outline,
        slide_spec_file=slide_spec,
        require_adaptive_visual_plan=True,
        source_inventory_entries=entries,
        semantic_gate=True,
    )
    assumed, outline_b, slide_spec_b = _semantic_fixture(
        tmp_path,
        body_point="假设：转化率有望提升 30%",
        source_refs=["src-01"],
    )
    validate_presentation_html_contract(
        assumed,
        expected_page_count=1,
        outline_file=outline_b,
        slide_spec_file=slide_spec_b,
        require_adaptive_visual_plan=True,
        source_inventory_entries=entries,
        semantic_gate=True,
    )


def test_semantic_gate_off_keeps_v1_contract_behavior(tmp_path: Path) -> None:
    # Without the v2 seam the same unsourced claim passes the v1 contract —
    # historical decks and quick conversions are untouched.
    source, outline, slide_spec = _semantic_fixture(
        tmp_path,
        body_point="市场占有率 48%",
        source_refs=["src-01"],
    )
    validate_presentation_html_contract(
        source,
        expected_page_count=1,
        outline_file=outline,
        slide_spec_file=slide_spec,
        require_adaptive_visual_plan=True,
    )


# ─── FR-P4 parameterized font / contrast gates ──────────────────


def _layout_item(text: str, **overrides) -> dict:
    item = {
        "kind": "text",
        "text": text,
        "x": 80,
        "y": 120,
        "w": 600,
        "h": 60,
        "style": {},
    }
    item.update(overrides)
    return item


def test_text_bounds_parameterized_font_floor() -> None:
    layout = {
        "slides": [
            {
                "width": 1280,
                "height": 720,
                "items": [_layout_item("正文内容", style={"fontSize": "12px"})],
            }
        ]
    }
    # v1 callers keep bounds-only behavior.
    validate_browser_slide_text_bounds(layout)
    with pytest.raises(ValueError, match="minimum is 16"):
        validate_browser_slide_text_bounds(layout, minimum_body_font_size_px=16)

    metadata_layout = {
        "slides": [
            {
                "width": 1280,
                "height": 720,
                "items": [
                    _layout_item(
                        "页脚",
                        style={"fontSize": "12px"},
                        textRole="metadata",
                    )
                ],
            }
        ]
    }
    validate_browser_slide_text_bounds(
        metadata_layout,
        minimum_body_font_size_px=16,
        minimum_metadata_font_size_px=10,
    )
    with pytest.raises(ValueError, match="minimum is 12"):
        validate_browser_slide_text_bounds(
            metadata_layout,
            minimum_metadata_font_size_px=12.5,
        )


def test_text_bounds_parameterized_contrast_floor() -> None:
    low_contrast = {
        "slides": [
            {
                "width": 1280,
                "height": 720,
                "backgroundColor": "#ffffff",
                "items": [
                    _layout_item(
                        "低对比度正文",
                        style={"color": "#777777", "fontSize": "18px"},
                    )
                ],
            }
        ]
    }
    validate_browser_slide_text_bounds(low_contrast)
    with pytest.raises(ValueError, match="contrast"):
        validate_browser_slide_text_bounds(low_contrast, minimum_contrast_ratio=4.5)

    high_contrast = {
        "slides": [
            {
                "width": 1280,
                "height": 720,
                "items": [
                    _layout_item(
                        "高对比度正文",
                        style={"color": "rgb(0, 0, 0)", "backgroundColor": "#ffffff"},
                    )
                ],
            }
        ]
    }
    validate_browser_slide_text_bounds(high_contrast, minimum_contrast_ratio=4.5)


def test_text_contrast_composites_css_alpha_against_slide_background() -> None:
    transparent_item_background = {
        "slides": [
            {
                "width": 1280,
                "height": 720,
                "backgroundColor": "rgb(255, 255, 255)",
                "items": [
                    _layout_item(
                        "透明元素背景不应被误判为黑色",
                        style={
                            "color": "rgb(19, 33, 58)",
                            "backgroundColor": "rgba(0, 0, 0, 0)",
                        },
                    )
                ],
            }
        ]
    }
    validate_browser_slide_text_bounds(
        transparent_item_background,
        minimum_contrast_ratio=4.5,
    )

    transparent_foreground = {
        "slides": [
            {
                "width": 1280,
                "height": 720,
                "backgroundColor": "rgb(255, 255, 255)",
                "items": [
                    _layout_item(
                        "透明文字仍必须被质量门禁拒绝",
                        style={
                            "color": "rgba(0, 0, 0, 0)",
                            "backgroundColor": "rgba(0, 0, 0, 0)",
                        },
                    )
                ],
            }
        ]
    }
    with pytest.raises(ValueError, match="contrast"):
        validate_browser_slide_text_bounds(
            transparent_foreground,
            minimum_contrast_ratio=4.5,
        )


def test_visual_quality_accepts_parameterized_floors() -> None:
    layout = {
        "slides": [
            {
                "width": 1280,
                "height": 720,
                "items": [
                    _layout_item(
                        "较大正文",
                        style={"fontSize": "18px"},
                        clientWidth=600,
                        clientHeight=60,
                        scrollWidth=600,
                        scrollHeight=60,
                    )
                ],
            }
        ]
    }
    # The v1 floor (16px) passes; a raised v2 floor fails the same layout.
    validate_browser_slide_visual_quality(layout)
    with pytest.raises(PresentationVisualQualityError):
        validate_browser_slide_visual_quality(layout, minimum_body_font_size_px=20)


# ─── FR-P4/P5/P8 artifact-level deck gates ──────────────────────


def _deck_pptx_bytes(*, text_per_slide: int = 200, shapes_per_slide: int = 6) -> bytes:
    from io import BytesIO

    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    for index in range(8):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(5))
        box.text = f"第 {index + 1  } 页 " + "内容" * max(1, text_per_slide // 2)
        for extra in range(shapes_per_slide - 1):
            del extra
            slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(2), Inches(0.5))
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _deck_pptx_with_full_page_image(
    *,
    data_slide_index: int = 2,
    editable_shapes_on_data_slide: int = 1,
) -> bytes:
    from io import BytesIO

    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    image_output = BytesIO()
    Image.new("RGB", (1600, 900), color=(28, 34, 46)).save(image_output, format="PNG")
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    for index in range(8):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        if index == data_slide_index:
            image_output.seek(0)
            slide.shapes.add_picture(
                image_output, 0, 0, presentation.slide_width, presentation.slide_height
            )
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(5))
        box.text = f"第 {index + 1} 页 " + "内容" * 120
        if index == data_slide_index:
            for extra in range(1, editable_shapes_on_data_slide):
                extra_box = slide.shapes.add_textbox(
                    Inches(0.5 + extra * 1.25),
                    Inches(6),
                    Inches(1),
                    Inches(0.5),
                )
                extra_box.text = f"指标 {extra}"
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _verified(artifact_type: str, data: bytes) -> _VerifiedArtifact:
    import hashlib

    return _VerifiedArtifact(
        artifact_type=artifact_type,
        workspace_path=f"workspace/deliverables/x/final.{artifact_type}",
        content_hash=hashlib.sha256(data).hexdigest(),
        size_bytes=len(data),
        tool_call_id="tool-1",
        data=data,
    )


def _slide_spec(*, data_slides: tuple[int, ...] = ()) -> dict:
    return {
        "visual_plan_version": "adaptive-v1",
        "slides": [
            {
                "slide_id": f"slide-{index:02d}",
                "headline": f"第 {index} 页",
                "body_points": ["内容"],
                "source_refs": [],
                "visual_kind": "editable_table" if index in data_slides else "editable_typography",
                "data_slide": index in data_slides,
            }
            for index in range(1, 9)
        ],
    }


def test_density_gate_rejects_sparse_and_overstuffed_decks() -> None:
    request = _request()
    sparse = _deck_pptx_bytes(text_per_slide=8)
    facts, invalid = _presentation_contract_facts(
        request,
        {"pptx": _verified("pptx", sparse)},
        deck_slide_spec=_slide_spec(),
    )
    assert "pptx" in invalid
    assert facts["pptx"]["density_gate"] == 0

    dense = _deck_pptx_bytes(text_per_slide=200)
    facts, invalid = _presentation_contract_facts(
        request,
        {"pptx": _verified("pptx", dense)},
        deck_slide_spec=_slide_spec(),
    )
    assert not invalid
    assert facts["pptx"]["density_gate"] == 1

    # v1 requests keep the historical contract: no density gate at all.
    v1_request = _request(workflow_id="builtin.presentation.v1", workflow_version="1.0.0")
    _, invalid = _presentation_contract_facts(
        v1_request,
        {"pptx": _verified("pptx", sparse)},
    )
    assert not invalid


def test_data_slide_full_page_raster_is_rejected_for_v2() -> None:
    request = _request()
    rasterized = _deck_pptx_with_full_page_image(data_slide_index=2)
    facts, invalid = _presentation_contract_facts(
        request,
        {"pptx": _verified("pptx", rasterized)},
        deck_slide_spec=_slide_spec(data_slides=(3,)),
    )
    assert "pptx" in invalid
    assert facts["pptx"]["data_slide_editability_gate"] == 0

    # The same rasterized page is legal when the slide is not a data slide.
    facts, invalid = _presentation_contract_facts(
        request,
        {"pptx": _verified("pptx", rasterized)},
        deck_slide_spec=_slide_spec(),
    )
    assert not invalid
    assert facts["pptx"]["data_slide_editability_gate"] == 1


def test_data_slide_allows_full_bleed_background_with_editable_composition() -> None:
    request = _request()
    hybrid = _deck_pptx_with_full_page_image(
        data_slide_index=2,
        editable_shapes_on_data_slide=4,
    )
    facts, invalid = _presentation_contract_facts(
        request,
        {"pptx": _verified("pptx", hybrid)},
        deck_slide_spec=_slide_spec(data_slides=(3,)),
    )

    assert not invalid
    assert facts["pptx"]["picture_coverage_ratio_by_slide"][2] == 1
    assert facts["pptx"]["editable_shape_count_by_slide"][2] == 4
    assert facts["pptx"]["data_slide_editability_gate"] == 1


def test_v2_deck_requires_slide_spec_evidence() -> None:
    request = _request()
    facts, invalid = _presentation_contract_facts(
        request,
        {"pptx": _verified("pptx", _deck_pptx_bytes())},
        deck_slide_spec=None,
    )
    assert "pptx" in invalid
    assert facts["pptx"]["slide_spec_gate"] == 0


def _pdf_pages_with_content(*, blank_pages: tuple[int, ...] = ()) -> bytes:
    import fitz

    document = fitz.open()
    for index in range(8):
        page = document.new_page(width=960, height=540)
        if index not in blank_pages:
            page.draw_rect(fitz.Rect(40, 40, 920, 500), color=(0.2, 0.3, 0.8), width=8)
            page.insert_text((72, 96), f"Slide {index + 1} 内容")
    data = document.tobytes()
    document.close()
    return data


def test_pdf_contract_boundary_and_consistency_spot_check() -> None:
    pptx_data = _deck_pptx_bytes()
    good_pdf = _pdf_pages_with_content()
    blank_pdf = _pdf_pages_with_content(blank_pages=(2,))

    # Default PPTX-only contract: a PDF is never required, so PDF rendering
    # issues can never block the delivery.
    pptx_only = _request(output_contract=["pptx"])
    facts, invalid = _presentation_contract_facts(
        pptx_only,
        {"pptx": _verified("pptx", pptx_data)},
        deck_slide_spec=_slide_spec(),
    )
    assert not invalid
    assert "pdf" not in facts

    # Explicit PDF contract: a blank rendered page where the PPTX slide has
    # real text fails the consistency spot check.
    with_pdf = _request(output_contract=["pptx", "pdf"])
    facts, invalid = _presentation_contract_facts(
        with_pdf,
        {"pptx": _verified("pptx", pptx_data), "pdf": _verified("pdf", blank_pdf)},
        deck_slide_spec=_slide_spec(),
    )
    assert {"pptx", "pdf"} <= set(invalid)
    assert facts["pdf"]["visual_consistency_gate"] == 0

    facts, invalid = _presentation_contract_facts(
        with_pdf,
        {"pptx": _verified("pptx", pptx_data), "pdf": _verified("pdf", good_pdf)},
        deck_slide_spec=_slide_spec(),
    )
    assert not invalid
    assert facts["pdf"]["visual_consistency_gate"] == 1

    # v1 keeps page-count parity only; the new spot check never applies.
    v1_request = _request(
        workflow_id="builtin.presentation.v1",
        workflow_version="1.0.0",
        output_contract=["pptx", "pdf"],
    )
    _, invalid = _presentation_contract_facts(
        v1_request,
        {"pptx": _verified("pptx", pptx_data), "pdf": _verified("pdf", blank_pdf)},
    )
    assert not invalid
