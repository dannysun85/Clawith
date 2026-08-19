"""FR-P6 page-targeted revision runtime + FR-P7 font substitution reports."""

from __future__ import annotations

from io import BytesIO
from unittest.mock import AsyncMock
import uuid

from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import pytest

from app.models.deliverable import (
    DeliverableExecution,
    DeliverableRequest,
)
from app.services.deliverable_artifacts import _pptx_facts
from app.services.deliverable_executions import execution_unit_blueprints
from app.services.deliverable_workflows import (
    DeliverableWorkflowError,
    build_deliverable_prompt,
    prepare_deliverable_launch,
)
from app.services.document_conversion.font_report import (
    available_font_families,
    font_substitution_report,
    requested_font_families_from_html,
    requested_font_families_from_pptx,
)


class _Result:
    def __init__(self, value: object | None) -> None:
        self.value = value

    def scalar_one_or_none(self):
        return self.value

    def scalars(self):
        return self

    def all(self):
        if self.value is None:
            return []
        return self.value if isinstance(self.value, list) else [self.value]


class _Session:
    def __init__(self, *execute_values: object | None) -> None:
        self.execute_values = list(execute_values)
        self.added: list[object] = []

    async def execute(self, _statement):
        return _Result(self.execute_values.pop(0))

    async def flush(self) -> None:
        pass

    async def get(self, _model, _key):
        return None

    def add(self, value: object) -> None:
        self.added.append(value)


def _spec(**overrides) -> dict:
    spec = {
        "audience": "渠道运营团队",
        "scenario": "季度复盘汇报",
        "page_count": 8,
        "language": "zh-CN",
        "style": "商务简洁",
        "editability_contract": "editable",
    }
    spec.update(overrides)
    return spec


def _request(**overrides) -> DeliverableRequest:
    request = DeliverableRequest(
        id=uuid.uuid4(),
        tenant_id=uuid.uuid4(),
        created_by_user_id=uuid.uuid4(),
        agent_id=uuid.uuid4(),
        session_id=uuid.uuid4(),
        client_request_id=uuid.uuid4(),
        request_fingerprint="a" * 64,
        work_type="presentation",
        workflow_id="builtin.presentation.v2",
        workflow_version="2.0.0",
        goal="为保温杯季度销售复盘制作汇报 PPT",
        inputs=[],
        spec=_spec(),
        tier="pro",
        approval_policy=["outline", "final"],
        output_contract=["pptx"],
        status="ready",
        current_stage="revision_ready",
        version=5,
        contract_revision=2,
    )
    for key, value in overrides.items():
        setattr(request, key, value)
    return request


def _revision_execution(
    request: DeliverableRequest,
    *,
    target_units: list[str],
    instruction: str = "第 3 页减少文字并突出核心数据",
) -> DeliverableExecution:
    execution = DeliverableExecution(
        id=uuid.uuid4(),
        tenant_id=request.tenant_id,
        request_id=request.id,
        execution_number=2,
        kind="revision",
        status="ready",
        current_stage="revision_ready",
        workflow_id=request.workflow_id,
        workflow_version=request.workflow_version,
        contract_snapshot={
            "version": 1,
            "request_id": str(request.id),
            "revision_instruction": instruction,
            "target_units": target_units,
        },
        preflight_snapshot={},
        revision_instruction=instruction,
        idempotency_key=uuid.uuid4(),
        request_fingerprint="b" * 64,
    )
    request.current_execution_id = execution.id
    return execution


# ─── FR-P6: blueprint keeps only target slides ──────────────────


def test_revision_blueprint_only_reruns_target_slides() -> None:
    request = _request()
    blueprints = execution_unit_blueprints(request, target_units=["slide-03"])
    slide_units = [bp for bp in blueprints if bp.stage_key == "slide_render"]
    assert [bp.unit_key for bp in slide_units] == ["slide-03"]
    # Deck-level assembly/QA stages always rerun; untouched slides never do.
    stage_keys = {bp.stage_key for bp in blueprints}
    assert {"deck_assemble", "pptx_render", "semantic_qa"} <= stage_keys
    assert not any(bp.unit_key == "slide-01" for bp in blueprints)
    full = execution_unit_blueprints(request)
    assert sum(1 for bp in full if bp.stage_key == "slide_render") == 8


def test_slide_revision_prompt_scopes_changes_to_target_slides() -> None:
    request = _request()
    prompt = build_deliverable_prompt(
        request,
        presentation_v2_stage="slide_revision",
        presentation_v2_target_units=["slide-03"],
        presentation_v2_revision_instruction="第 3 页减少文字并突出核心数据",
    )
    assert "slide-03" in prompt
    assert "第 3 页减少文字并突出核心数据" in prompt
    assert "render_mode='hybrid_editable'" in prompt
    assert "convert_html_to_pptx" in prompt
    assert "convert_html_to_pdf" not in prompt  # PDF only when contracted
    assert "never add a new unsourced number" in prompt

    pdf_request = _request(output_contract=["pptx", "pdf"])
    pdf_prompt = build_deliverable_prompt(
        pdf_request,
        presentation_v2_stage="slide_revision",
        presentation_v2_target_units=["slide-02", "slide-05"],
        presentation_v2_revision_instruction="更换第 2 页主视觉并压缩第 5 页表格",
    )
    assert "slide-02, slide-05" in pdf_prompt
    assert "convert_html_to_pdf" in pdf_prompt

    with pytest.raises(DeliverableWorkflowError):
        build_deliverable_prompt(
            request,
            presentation_v2_stage="slide_revision",
            presentation_v2_target_units=[],
            presentation_v2_revision_instruction="x",
        )


@pytest.mark.asyncio
async def test_revision_launch_enters_slide_revision_without_redrafting(monkeypatch) -> None:
    request = _request()
    execution = _revision_execution(request, target_units=["slide-03"])
    monkeypatch.setattr(
        "app.services.deliverable_workflows.preflight_workflow",
        AsyncMock(return_value={"launchable": True, "reasons": []}),
    )
    db = _Session(
        request,
        execution,
        None,  # brief row lookup for the inventory projection
    )
    prepared = await prepare_deliverable_launch(
        db,  # type: ignore[arg-type]
        request_id=request.id,
        tenant_id=request.tenant_id,
        user_id=request.created_by_user_id,
        agent_id=request.agent_id,
        session_id=request.session_id,
        message_id=uuid.uuid4(),
    )
    assert request.status == "running"
    assert request.current_stage == "slide_revision"
    assert "slide-03" in prepared.prompt
    assert "outline.json" in prepared.prompt  # reads back, not re-drafted
    assert "only these slides are in scope" in prepared.prompt


@pytest.mark.asyncio
async def test_full_deck_revision_without_targets_redrafts_outline(monkeypatch) -> None:
    request = _request()
    execution = _revision_execution(request, target_units=[])
    monkeypatch.setattr(
        "app.services.deliverable_workflows.preflight_workflow",
        AsyncMock(return_value={"launchable": True, "reasons": []}),
    )
    db = _Session(request, execution, None)
    prepared = await prepare_deliverable_launch(
        db,  # type: ignore[arg-type]
        request_id=request.id,
        tenant_id=request.tenant_id,
        user_id=request.created_by_user_id,
        agent_id=request.agent_id,
        session_id=request.session_id,
        message_id=uuid.uuid4(),
    )
    assert request.current_stage == "outline_draft"
    assert "slide-03" not in prepared.prompt


# ─── FR-P7: font substitution report ────────────────────────────


def test_requested_families_from_html_and_pure_report() -> None:
    html = (
        "<style>.t { font-family: 'MiSans', sans-serif; }</style>"
        '<p style="font-family: &quot;Noto Sans SC&quot;, serif">x</p>'
        '<span style="font-family: monospace">y</span>'
    )
    requested = requested_font_families_from_html(html)
    assert requested == ("misans", "sans-serif", "noto sans sc", "serif", "monospace")
    report = font_substitution_report(requested, frozenset({"noto sans sc"}))
    # monospace is a generic alias; only MiSans is missing on this host.
    assert report == [
        {
            "requested": "misans",
            "actual": "host default sans/serif fallback",
            "reason": "font_not_installed",
        }
    ]


def _pptx_with_fonts(families: list[str]) -> bytes:
    image_output = BytesIO()
    Image.new("RGB", (1600, 900), color=(28, 34, 46)).save(image_output, format="PNG")
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    for index, family in enumerate(families):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = f"第 {index + 1} 页正文内容，用于字体报告测试。"
        run.font.name = family
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def test_pptx_facts_record_font_substitutions() -> None:
    data = _pptx_with_fonts(["Astra Definitely Missing Font 2049", "Arial"])
    facts = _pptx_facts(data)
    assert "astra definitely missing font 2049" in facts["font_families_requested"]
    assert facts["font_substitutions"] == [
        {
            "requested": "astra definitely missing font 2049",
            "actual": "host default sans/serif fallback",
            "reason": "font_not_installed",
        }
    ]


def test_pptx_facts_report_no_substitution_for_installed_fonts() -> None:
    data = _pptx_with_fonts(["Arial"])
    facts = _pptx_facts(data)
    assert facts["font_substitutions"] == []


def test_requested_families_from_pptx_reads_typefaces() -> None:
    data = _pptx_with_fonts(["MiSans"])
    assert "misans" in requested_font_families_from_pptx(data)


def test_available_font_discovery_is_cached_and_non_empty_or_empty() -> None:
    families = available_font_families()
    assert isinstance(families, frozenset)
    again = available_font_families()
    assert again is families
