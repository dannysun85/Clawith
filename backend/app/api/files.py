"""File management API routes for agent workspaces."""

import asyncio
import base64
import csv
import io
import mimetypes
import uuid
from pathlib import Path, PurePosixPath

import aiofiles
from fastapi import APIRouter, Depends, File as FastFile, HTTPException, Request, status
from fastapi import UploadFile as UploadFileType
from fastapi.responses import FileResponse, Response
from fastapi.security import HTTPAuthorizationCredentials, HTTPBearer
from pydantic import BaseModel

from app.config import get_settings
from app.core.permissions import check_agent_access
from app.core.security import (
    BROWSER_SESSION_COOKIE,
    get_current_user,
    get_current_user_for_access_token,
)
from app.database import get_db
from app.models.user import User
from app.models.workspace import WorkspaceFileRevision
from app.services.focus_service import is_focus_file_path
from app.services.workspace_collaboration import (
    acquire_edit_lock,
    content_hash,
    delete_workspace_file,
    list_revisions,
    read_text_if_exists,
    release_edit_lock,
    write_workspace_file,
)
from app.services.storage import (
    ensure_local_path,
    get_storage_backend,
    guess_content_type,
    normalize_storage_key,
)
from app.services.storage_runtime.base import StorageEntry
from app.services.workspace_paths import WorkspacePathError, resolve_agent_visible_path
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.orm import selectinload

settings = get_settings()
router = APIRouter(prefix="/agents/{agent_id}/files", tags=["files"])

_BINARY_CONTENT_SUFFIXES = {
    ".7z", ".avi", ".bmp", ".doc", ".docx", ".flac", ".gif", ".gz",
    ".jpeg", ".jpg", ".m4a", ".mkv", ".mov", ".mp3", ".mp4", ".ogg",
    ".pdf", ".png", ".ppt", ".pptx", ".rar", ".tar", ".wav", ".webm",
    ".webp", ".xls", ".xlsx", ".zip",
}

_INTERNAL_SKILL_METADATA_FILES = frozenset({
    ".astra-import.json",
    ".astra-managed.json",
})


def _is_internal_skill_metadata_path(path: str) -> bool:
    """Return whether a workspace path points at product-owned Skill metadata."""
    normalized = (path or "").strip().strip("/")
    parts = PurePosixPath(normalized).parts
    return (
        len(parts) >= 3
        and parts[0] == "skills"
        and parts[-1] in _INTERNAL_SKILL_METADATA_FILES
    )


def _reject_internal_skill_metadata_path(path: str, *, mutation: bool = False) -> None:
    """Keep Skill provenance manifests outside the user-editable file surface."""
    if not _is_internal_skill_metadata_path(path):
        return
    if mutation:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Managed Skill metadata cannot be modified directly",
        )
    raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found")


async def _read_text_or_binary_placeholder(storage, key: str, path: str) -> str:
    """Read strict UTF-8 text without ever decoding known binary media."""
    if Path(path).suffix.lower() in _BINARY_CONTENT_SUFFIXES:
        stat = await storage.stat(key)
        return f"[二进制文件: {Path(path).name}, {stat.size} bytes]"
    try:
        return await storage.read_text(key, encoding="utf-8", errors="strict")
    except UnicodeDecodeError:
        stat = await storage.stat(key)
        return f"[二进制文件: {Path(path).name}, {stat.size} bytes]"


class FileInfo(BaseModel):
    name: str
    path: str
    is_dir: bool
    size: int = 0
    modified_at: str = ""
    version_token: str | None = None
    url: str | None = None


class FileContent(BaseModel):
    path: str
    content: str
    version_token: str | None = None


class FileWrite(BaseModel):
    content: str
    autosave: bool = False
    session_id: str | None = None
    expected_version_token: str | None = None


class FileLockBody(BaseModel):
    path: str
    session_id: str | None = None


async def _directory_total_size(storage, storage_key: str) -> int:
    """Return the recursive byte size of all files below a storage directory."""
    total = 0
    pending = [normalize_storage_key(storage_key)]
    visited: set[str] = set()
    while pending:
        current = pending.pop()
        if current in visited:
            continue
        visited.add(current)
        for entry in await storage.list_dir(current):
            if entry.is_dir:
                pending.append(normalize_storage_key(entry.key))
            else:
                total += max(0, entry.size)
    return total


class RestoreRevisionBody(BaseModel):
    revision_id: uuid.UUID
    expected_version_token: str | None = None


TEXT_PREVIEW_EXTENSIONS = {
    ".bat",
    ".bash",
    ".c",
    ".cfg",
    ".clj",
    ".cpp",
    ".cs",
    ".css",
    ".dart",
    ".env",
    ".go",
    ".h",
    ".hpp",
    ".ini",
    ".java",
    ".js",
    ".jsx",
    ".kt",
    ".kts",
    ".less",
    ".lua",
    ".m",
    ".mm",
    ".php",
    ".pl",
    ".pm",
    ".properties",
    ".py",
    ".r",
    ".rb",
    ".rs",
    ".sass",
    ".scala",
    ".scss",
    ".sh",
    ".sql",
    ".swift",
    ".toml",
    ".ts",
    ".tsx",
    ".vue",
    ".xml",
    ".yaml",
    ".yml",
    ".zsh",
}

TEXT_PREVIEW_FILENAMES = {
    ".dockerignore",
    ".env",
    ".env.example",
    ".gitignore",
    ".npmrc",
    ".prettierrc",
    "dockerfile",
    "makefile",
}


def _agent_base_dir(agent_id: uuid.UUID) -> Path:
    local_root = settings.STORAGE_LOCAL_ROOT or settings.AGENT_DATA_DIR
    primary = Path(local_root) / str(agent_id)
    legacy = Path(settings.AGENT_DATA_DIR) / str(agent_id)
    if primary != legacy and legacy.exists() and not primary.exists():
        return legacy
    return primary


def _agent_storage_key(agent_id: uuid.UUID, rel_path: str = "") -> str:
    prefix = str(agent_id)
    rel = normalize_storage_key(rel_path)
    return f"{prefix}/{rel}" if rel else prefix


def _safe_path(agent_id: uuid.UUID, rel_path: str) -> Path:
    """Ensure the path is within the agent's directory (no path traversal)."""
    base = _agent_base_dir(agent_id)
    full = (base / rel_path).resolve()
    if not str(full).startswith(str(base.resolve())):
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Path traversal not allowed")
    return full


def _visible_path(agent_id: uuid.UUID, rel_path: str, tenant_id: uuid.UUID | None) -> tuple[Path, Path, bool]:
    """Resolve an agent-visible path, including virtual enterprise_info/."""
    try:
        resolved = resolve_agent_visible_path(
            _agent_base_dir(agent_id),
            rel_path,
            workspace_root=Path(settings.AGENT_DATA_DIR),
            tenant_id=str(tenant_id) if tenant_id else None,
        )
    except WorkspacePathError as exc:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail=str(exc)) from exc
    return resolved.path, resolved.relative_root, resolved.is_enterprise


def _is_enterprise_visible_path(rel_path: str) -> bool:
    normalized = (rel_path or "").strip().strip("/")
    return normalized == "enterprise_info" or normalized.startswith("enterprise_info/")


def _visible_storage_key(agent_id: uuid.UUID, rel_path: str, tenant_id: uuid.UUID | None) -> tuple[str, bool]:
    normalized = (rel_path or "").strip().strip("/")
    if _is_enterprise_visible_path(normalized):
        if not tenant_id:
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="No tenant associated")
        sub_path = normalized[len("enterprise_info"):].lstrip("/")
        return _enterprise_storage_key(str(tenant_id), sub_path), True
    return _agent_storage_key(agent_id, normalized), False


async def _require_agent_file_delete_access(
    db: AsyncSession,
    current_user: User,
    agent_id: uuid.UUID,
) -> None:
    """Allow destructive workspace file operations only for managers/admins."""
    _agent, access_level = await check_agent_access(db, current_user, agent_id)
    if access_level == "manage" or current_user.role in ("platform_admin", "org_admin", "super_admin"):
        return
    raise HTTPException(
        status_code=status.HTTP_403_FORBIDDEN,
        detail="Only agent managers or admins can delete files",
    )


@router.get("/", response_model=list[FileInfo])
async def list_files(
    agent_id: uuid.UUID,
    path: str = "",
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """List files and directories in an agent's file system."""
    await check_agent_access(db, current_user, agent_id)
    storage = get_storage_backend()
    storage_key, is_enterprise = _visible_storage_key(agent_id, path, current_user.tenant_id)
    normalized_path = (path or "").strip().strip("/")
    path_exists = await storage.exists(storage_key)
    path_is_dir = await storage.is_dir(storage_key)
    if not path_exists and not path_is_dir:
        if not (
            normalized_path in {"", "workspace"}
            or (is_enterprise and normalized_path == "enterprise_info")
        ):
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Path not found")
    elif path_exists and not path_is_dir:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Path is not a directory")

    items = []
    if not path and current_user.tenant_id:
        items.append(FileInfo(
            name="enterprise_info",
            path="enterprise_info",
            is_dir=True,
            size=0,
            modified_at="",
            version_token=None,
            url=None,
        ))
    entries = await storage.list_dir(storage_key) if path_exists or path_is_dir else []
    is_skills_path = normalized_path == "skills" or normalized_path.startswith("skills/")
    directory_entries = [entry for entry in entries if entry.is_dir] if is_skills_path else []
    directory_sizes = dict(zip(
        (entry.key for entry in directory_entries),
        await asyncio.gather(*(_directory_total_size(storage, entry.key) for entry in directory_entries)),
        strict=True,
    ))
    for entry in entries:
        if entry.name == '.gitkeep':
            continue
        if is_skills_path and entry.name in _INTERNAL_SKILL_METADATA_FILES:
            continue
        # Focus is database-backed. Legacy focus.md must stay hidden so users
        # do not edit a stale second source of truth.
        if not path and entry.name.lower() == "focus.md":
            continue
        if not path and entry.name.lower() == "agenda.md":
            continue
        if not path and entry.name == "enterprise_info":
            continue
        if is_enterprise:
            rel = str(Path(entry.key).relative_to(f"enterprise_info_{current_user.tenant_id}"))
            rel_path = f"enterprise_info/{rel}" if rel != "." else "enterprise_info"
        else:
            rel_path = str(Path(entry.key).relative_to(str(agent_id)))
        version_token = _entry_version_token(entry)
        if not entry.is_dir and not (entry.version_id or entry.etag or entry.content_hash):
            version = await storage.get_version(entry.key)
            version_token = version.token or version_token
        items.append(FileInfo(
            name=entry.name,
            path=rel_path,
            is_dir=entry.is_dir,
            size=directory_sizes.get(entry.key, entry.size),
            modified_at=entry.modified_at,
            version_token=version_token,
            url=f"/api/agents/{agent_id}/files/download?path={rel_path}" if not entry.is_dir else None
        ))
    return items


@router.get("/content", response_model=FileContent)
async def read_file(
    agent_id: uuid.UUID,
    path: str,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Read the content of a file."""
    _reject_internal_skill_metadata_path(path)
    storage = get_storage_backend()
    key, _ = _visible_storage_key(agent_id, path, current_user.tenant_id)
    await check_agent_access(db, current_user, agent_id)
    if is_focus_file_path(path):
        raise HTTPException(
            status_code=status.HTTP_410_GONE,
            detail="Focus is stored in the system database. Use the Focus API.",
        )
    if not await storage.exists(key) or not await storage.is_file(key):
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found")
    version = await storage.get_version(key)

    content = await _read_text_or_binary_placeholder(storage, key, path)
    return FileContent(path=path, content=content, version_token=version.token)


def _entry_version_token(entry: StorageEntry) -> str | None:
    token = entry.version_id or entry.etag or entry.content_hash
    if token:
        return token
    if entry.is_dir:
        return None
    if entry.modified_at or entry.size:
        return f"{entry.modified_at}:{entry.size}"
    return None


def _file_kind(path: str) -> str:
    file_path = Path(path)
    ext = file_path.suffix.lower()
    name = file_path.name.lower()
    if ext in {".md", ".markdown"}:
        return "markdown"
    if ext == ".csv":
        return "csv"
    if ext in {".html", ".htm"}:
        return "html"
    if ext == ".pdf":
        return "pdf"
    if ext in {".xlsx", ".xls"}:
        return "xlsx"
    if ext in {".docx", ".doc"}:
        return "docx"
    if ext in {".pptx", ".ppt"}:
        return "pptx"
    if ext in {".txt", ".log", ".json"} or ext in TEXT_PREVIEW_EXTENSIONS or name in TEXT_PREVIEW_FILENAMES:
        return "text"
    if ext in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg"}:
        return "image"
    if ext in {".mp4", ".webm", ".mov", ".m4v"}:
        return "video"
    if ext in {".mp3", ".wav", ".ogg", ".m4a", ".aac", ".flac", ".opus"}:
        return "audio"
    return "binary"


def _find_companion_text_preview(target: Path) -> Path | None:
    for suffix in (".md", ".txt"):
        candidate = target.with_suffix(suffix)
        if candidate.exists() and candidate.is_file():
            return candidate
    return None


def _extract_document_text(target: Path, kind: str) -> str:
    """Best-effort rich document text extraction for lightweight previews."""
    try:
        if kind == "xlsx":
            from openpyxl import load_workbook

            wb = load_workbook(target, read_only=True, data_only=True)
            sheets: list[str] = []
            for ws in wb.worksheets[:5]:
                rows = []
                for row in ws.iter_rows(max_row=80, max_col=20, values_only=True):
                    rows.append("\t".join("" if cell is None else str(cell) for cell in row))
                sheets.append(f"Sheet: {ws.title}\n" + "\n".join(rows))
            return "\n\n".join(sheets)
        if kind == "docx":
            from docx import Document

            doc = Document(str(target))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if kind == "pptx":
            from pptx import Presentation

            prs = Presentation(str(target))
            slides = []
            for idx, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                slides.append(f"Slide {idx}\n" + "\n".join(texts))
            return "\n\n".join(slides)
    except ImportError as exc:
        return f"Missing preview dependency: {exc}"
    except Exception as exc:
        return f"Preview extraction failed: {str(exc)[:200]}"
    return ""


def _detect_csv_delimiter(text: str) -> str:
    lines = [line.strip() for line in text.splitlines() if line.strip()][:10]
    if not lines:
        return ","
    candidates = [",", "，", ";", "\t", "|"]
    scores = {
        candidate: sum(line.count(candidate) for line in lines)
        for candidate in candidates
    }
    return max(scores, key=scores.get) if any(scores.values()) else ","


def _parse_csv_rows(text: str) -> list[list[str]]:
    delimiter = _detect_csv_delimiter(text)
    rows = list(csv.reader(io.StringIO(text), delimiter=delimiter))
    normalized: list[list[str]] = []
    for row in rows[:500]:
        values = list(row)
        while values and not str(values[-1] or "").strip():
            values.pop()
        if values:
            normalized.append(values)
    return normalized


@router.get("/preview")
async def preview_file(
    agent_id: uuid.UUID,
    path: str,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Return a browser-friendly preview payload for Workspace files."""
    _reject_internal_skill_metadata_path(path)
    await check_agent_access(db, current_user, agent_id)
    storage = get_storage_backend()
    key, _ = _visible_storage_key(agent_id, path, current_user.tenant_id)
    if not await storage.exists(key) or not await storage.is_file(key):
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found")

    kind = _file_kind(path)
    mime_type = mimetypes.guess_type(Path(path).name)[0] or "application/octet-stream"
    download_url = f"/api/agents/{agent_id}/files/download?path={path}"
    local_target: Path | None = None

    if kind in {"markdown", "html", "text"}:
        content = await storage.read_text(key, encoding="utf-8", errors="replace")
        return {
            "path": path,
            "kind": kind,
            "mime_type": mime_type,
            "content": content or "",
            "content_hash": content_hash(content or ""),
            "download_url": download_url,
        }
    if kind == "csv":
        content = await storage.read_text(key, encoding="utf-8", errors="replace")
        rows = _parse_csv_rows(content)
        return {
            "path": path,
            "kind": kind,
            "mime_type": mime_type,
            "content": content,
            "content_hash": content_hash(content),
            "rows": rows[:500],
            "download_url": download_url,
        }
    if kind in {"image", "pdf", "video", "audio"}:
        return {
            "path": path,
            "kind": kind,
            "mime_type": mime_type,
            "url": download_url,
            "download_url": download_url,
        }
    if kind == "xlsx":
        try:
            target = await ensure_local_path(key)
            local_target = target
            from openpyxl import load_workbook

            wb = load_workbook(target, read_only=True, data_only=True)
            sheets = []
            for ws in wb.worksheets[:5]:
                rows = []
                for row in ws.iter_rows(max_row=120, max_col=30, values_only=True):
                    values = ["" if cell is None else str(cell) for cell in row]
                    while values and not str(values[-1] or "").strip():
                        values.pop()
                    if any(value.strip() for value in values):
                        rows.append(values)
                sheets.append({
                    "title": ws.title,
                    "rows": rows,
                })
            wb.close()
            return {
                "path": path,
                "kind": kind,
                "mime_type": mime_type,
                "text": _extract_document_text(target, kind),
                "sheets": sheets,
                "download_url": download_url,
            }
        except Exception as exc:
            return {
                "path": path,
                "kind": kind,
                "mime_type": mime_type,
                "text": f"Preview extraction failed: {str(exc)[:200]}",
                "download_url": download_url,
            }
    if kind in {"docx", "pptx"}:
        target = await ensure_local_path(key)
        local_target = target
        extracted_text = _extract_document_text(target, kind)
        companion = _find_companion_text_preview(target)
        companion_content = await read_text_if_exists(companion) if companion is not None else None
        return {
            "path": path,
            "kind": kind,
            "mime_type": mime_type,
            "text": companion_content or extracted_text,
            "companion_path": str(companion.resolve().relative_to(_agent_base_dir(agent_id).resolve())) if companion is not None and not path.startswith("enterprise_info") else None,
            "download_url": download_url,
        }

    if local_target is not None:
        companion = _find_companion_text_preview(local_target)
    else:
        companion = None
    if companion is not None:
        content = await read_text_if_exists(companion)
        return {
            "path": path,
            "kind": "text",
            "mime_type": "text/markdown" if companion.suffix.lower() == ".md" else "text/plain",
            "content": content or "",
            "content_hash": content_hash(content or ""),
            "companion_path": str(companion.resolve().relative_to(_agent_base_dir(agent_id).resolve())) if not path.startswith("enterprise_info") else None,
            "download_url": download_url,
        }

    raw = await storage.read_bytes(key)
    encoded = base64.b64encode(raw[:1024 * 1024]).decode("ascii")
    return {
        "path": path,
        "kind": kind,
        "mime_type": mime_type,
        "size": len(raw),
        "base64_sample": encoded,
        "download_url": download_url,
    }


@router.get("/download")
async def download_file(
    agent_id: uuid.UUID,
    path: str,
    request: Request,
    token: str = "",
    inline: bool = False,
    credentials: HTTPAuthorizationCredentials | None = Depends(HTTPBearer(auto_error=False)),
    db: AsyncSession = Depends(get_db),
):
    """Download / serve a file from the agent workspace (browser-friendly).
    
    Auth via Bearer header or the same-origin HttpOnly browser session cookie.

    The query token remains a temporary compatibility fallback for cached
    frontend bundles and must not be generated by current clients.
    """
    _reject_internal_skill_metadata_path(path)

    # Resolve JWT token from Bearer, HttpOnly cookie, or legacy query fallback.
    jwt_token = None
    if credentials:
        jwt_token = credentials.credentials
    elif request.cookies.get(BROWSER_SESSION_COOKIE):
        jwt_token = request.cookies[BROWSER_SESSION_COOKIE]
    elif token:
        jwt_token = token

    if not jwt_token:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Authentication required")

    user = await get_current_user_for_access_token(jwt_token, db)

    await check_agent_access(db, user, agent_id)
    storage = get_storage_backend()
    key, _ = _visible_storage_key(agent_id, path, user.tenant_id)
    if not await storage.exists(key) or not await storage.is_file(key):
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found")
    presigned = await storage.presign_download_url(key, filename=Path(path).name, inline=inline)
    if presigned:
        return Response(
            status_code=302,
            headers={"Location": presigned},
        )
    local_path = await storage.local_path_for(key)
    if local_path is not None:
        return FileResponse(
            path=str(local_path),
            filename=Path(path).name,
            content_disposition_type="inline" if inline else "attachment",
        )
    data = await storage.read_bytes(key)
    disposition = "inline" if inline else "attachment"
    return Response(
        content=data,
        media_type=guess_content_type(Path(path).name),
        headers={"Content-Disposition": f'{disposition}; filename="{Path(path).name}"'},
    )


@router.put("/content")
async def write_file(
    agent_id: uuid.UUID,
    path: str,
    data: FileWrite,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Write content to a file (create or overwrite)."""
    await check_agent_access(db, current_user, agent_id)
    _reject_internal_skill_metadata_path(path, mutation=True)
    if is_focus_file_path(path):
        raise HTTPException(
            status_code=status.HTTP_410_GONE,
            detail="Focus is stored in the system database. Use the Focus API.",
        )
    if path.startswith("enterprise_info"):
        if current_user.role not in ("platform_admin", "org_admin"):
            raise HTTPException(status_code=403, detail="Only admins can edit enterprise knowledge base")
        if path.strip("/") == "enterprise_info":
            raise HTTPException(status_code=400, detail="Cannot overwrite enterprise_info root")
        target, _, _ = _visible_path(agent_id, path, current_user.tenant_id)
        target.parent.mkdir(parents=True, exist_ok=True)
        async with aiofiles.open(target, "w", encoding="utf-8") as f:
            await f.write(data.content)
        return {"status": "ok", "path": path, "revision_id": None}

    result = await write_workspace_file(
        db,
        agent_id=agent_id,
        base_dir=_agent_base_dir(agent_id),
        path=path,
        content=data.content,
        actor_type="user",
        actor_id=current_user.id,
        operation="autosave" if data.autosave else "write",
        session_id=data.session_id,
        enforce_human_lock=False,
        merge_user_autosave=data.autosave,
        expected_version_token=data.expected_version_token,
    )
    if not result.ok:
        raise HTTPException(status_code=status.HTTP_409_CONFLICT, detail=result.message)
    await db.commit()
    return {"status": "ok", "path": result.path, "revision_id": result.revision_id}


@router.post("/locks")
async def lock_file(
    agent_id: uuid.UUID,
    data: FileLockBody,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Acquire or refresh a short-lived human editing lock for a file."""
    await check_agent_access(db, current_user, agent_id)
    _reject_internal_skill_metadata_path(data.path, mutation=True)
    if is_focus_file_path(data.path):
        raise HTTPException(status_code=status.HTTP_410_GONE, detail="Focus is stored in the system database.")
    lock = await acquire_edit_lock(
        db,
        agent_id=agent_id,
        path=data.path,
        user_id=current_user.id,
        session_id=data.session_id,
    )
    await db.commit()
    return {"status": "ok", "path": lock.path, "expires_at": lock.expires_at.isoformat()}


@router.delete("/locks")
async def unlock_file(
    agent_id: uuid.UUID,
    path: str,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Release the current user's edit lock for a file."""
    await check_agent_access(db, current_user, agent_id)
    _reject_internal_skill_metadata_path(path, mutation=True)
    await release_edit_lock(db, agent_id=agent_id, path=path, user_id=current_user.id)
    if hasattr(db, "commit"):
        await db.commit()
    return {"status": "ok", "path": path}


@router.get("/revisions")
async def get_file_revisions(
    agent_id: uuid.UUID,
    path: str,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """List version history for the currently opened Workspace file."""
    await check_agent_access(db, current_user, agent_id)
    _reject_internal_skill_metadata_path(path)
    if is_focus_file_path(path):
        return []
    if path.startswith("enterprise_info"):
        return []
    revisions = await list_revisions(db, agent_id=agent_id, path=path)
    return [
        {
            "id": str(rev.id),
            "path": rev.path,
            "operation": rev.operation,
            "actor_type": rev.actor_type,
            "actor_id": str(rev.actor_id) if rev.actor_id else None,
            "session_id": rev.session_id,
            "before_content": rev.before_content,
            "after_content": rev.after_content,
            "created_at": rev.created_at.isoformat() if rev.created_at else None,
            "updated_at": rev.updated_at.isoformat() if rev.updated_at else None,
        }
        for rev in revisions
    ]


@router.post("/restore")
async def restore_file_revision(
    agent_id: uuid.UUID,
    data: RestoreRevisionBody,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Restore a file to a previous revision's after-content."""
    await check_agent_access(db, current_user, agent_id)
    result = await db.execute(
        select(WorkspaceFileRevision).where(
            WorkspaceFileRevision.id == data.revision_id,
            WorkspaceFileRevision.agent_id == agent_id,
        )
    )
    revision = result.scalar_one_or_none()
    if not revision:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Revision not found")
    _reject_internal_skill_metadata_path(revision.path, mutation=True)
    if revision.after_content is None:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Cannot restore an empty/deleted revision")

    restored = await write_workspace_file(
        db,
        agent_id=agent_id,
        base_dir=_agent_base_dir(agent_id),
        path=revision.path,
        content=revision.after_content,
        actor_type="user",
        actor_id=current_user.id,
        operation="restore",
        enforce_human_lock=False,
        expected_version_token=data.expected_version_token,
    )
    if not restored.ok:
        raise HTTPException(status_code=status.HTTP_409_CONFLICT, detail=restored.message)
    await db.commit()
    return {"status": "ok", "path": revision.path, "revision_id": restored.revision_id}


@router.delete("/content")
async def delete_file(
    agent_id: uuid.UUID,
    path: str,
    expected_version_token: str | None = None,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Delete a file."""
    await _require_agent_file_delete_access(db, current_user, agent_id)
    _reject_internal_skill_metadata_path(path, mutation=True)
    if is_focus_file_path(path):
        raise HTTPException(
            status_code=status.HTTP_410_GONE,
            detail="Focus is stored in the system database. Use the Focus API.",
        )
    if path.startswith("enterprise_info") and current_user.role not in ("platform_admin", "org_admin"):
        raise HTTPException(status_code=403, detail="Only admins can delete enterprise knowledge base files")
    if path.strip("/") == "enterprise_info":
        raise HTTPException(status_code=400, detail="Cannot delete enterprise_info root")
    result = await delete_workspace_file(
        db,
        agent_id=agent_id,
        base_dir=_agent_base_dir(agent_id),
        path=path,
        actor_type="user",
        actor_id=current_user.id,
        enforce_human_lock=False,
        expected_version_token=expected_version_token,
    )
    if not result.ok:
        if "not found" in result.message.lower():
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=result.message)
        raise HTTPException(status_code=status.HTTP_409_CONFLICT, detail=result.message)
    if hasattr(db, "commit"):
        await db.commit()
    return {"status": "ok", "path": path}


class ImportSkillBody(BaseModel):
    skill_id: uuid.UUID


@router.post("/import-skill")
async def import_skill_to_agent(
    agent_id: uuid.UUID,
    body: ImportSkillBody,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Import a global skill into this agent's skills/ workspace folder.

    Copies all files from the global skill registry into
    <agent_workspace>/skills/<folder_name>/.
    """
    agent, _ = await check_agent_access(db, current_user, agent_id)

    from app.models.skill import Skill
    from app.services.skill_scope import scope_skill_query

    # Load only a global or same-tenant skill visible to the target agent.
    result = await db.execute(
        scope_skill_query(
            select(Skill)
            .where(Skill.id == body.skill_id)
            .options(selectinload(Skill.files)),
            agent.tenant_id,
        )
    )
    skill = result.scalar_one_or_none()
    if not skill:
        raise HTTPException(status_code=404, detail="Skill not found")

    if not skill.files:
        raise HTTPException(status_code=400, detail="Skill has no files")

    from app.services.skill_workspace import deploy_skills_to_agent_workspace

    stats = await deploy_skills_to_agent_workspace(
        agent_id,
        [skill],
        provisioning="user_selected",
    )
    if stats["conflicts"]:
        raise HTTPException(
            status_code=status.HTTP_409_CONFLICT,
            detail=(
                "This Skill folder contains user changes and was left untouched. "
                "Rename or delete it before importing the registry version."
            ),
        )

    return {
        "status": "ok",
        "skill_name": skill.name,
        "folder_name": skill.folder_name,
        "files_written": stats["files"],
        "files": sorted(f.path for f in skill.files),
        "sync": stats,
    }


upload_router = APIRouter(prefix="/agents/{agent_id}/files", tags=["files"])
DEFAULT_UPLOAD_DIR = "workspace/uploads"


@upload_router.post("/upload")
async def upload_file_to_workspace(
    agent_id: uuid.UUID,
    file: UploadFileType = FastFile(...),
    path: str = "workspace/knowledge_base",
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Upload a binary file to agent workspace."""
    await check_agent_access(db, current_user, agent_id)

    normalized_path = (path or "").strip().strip("/")
    if not normalized_path or normalized_path == ".":
        normalized_path = DEFAULT_UPLOAD_DIR

    # Validate path prefix
    if normalized_path not in {"workspace", "skills"} and not normalized_path.startswith(("workspace/", "skills/")):
        raise HTTPException(status_code=400, detail="右侧根目录视图是 agent 根目录；上传文件时请放到 workspace/ 或 skills/ 目录下")

    filename = file.filename or "unnamed"
    # Sanitize filename
    filename = filename.replace("/", "_").replace("\\", "_")
    _reject_internal_skill_metadata_path(
        f"{normalized_path}/{filename}",
        mutation=True,
    )
    storage = get_storage_backend()
    file_key = _agent_storage_key(agent_id, f"{normalized_path}/{filename}")

    content = await file.read()
    await storage.write_bytes(file_key, content, content_type=guess_content_type(filename))

    # Auto-extract text from non-text files
    extracted_path = None
    from app.services.text_extractor import needs_extraction, save_extracted_text
    if needs_extraction(filename):
        save_path = await ensure_local_path(file_key)
        txt_file = save_extracted_text(save_path, content, filename)
        if txt_file:
            extracted_path = f"{normalized_path}/{txt_file.name}"
            extracted_key = _agent_storage_key(agent_id, extracted_path)
            await storage.write_bytes(extracted_key, txt_file.read_bytes(), content_type="text/plain; charset=utf-8")

    return {
        "status": "ok",
        "path": f"{normalized_path}/{filename}",
        "url": f"/api/agents/{agent_id}/files/download?path={normalized_path}/{filename}",
        "filename": filename,
        "size": len(content),
        "extracted_text_path": extracted_path,
    }


# ─── Enterprise Knowledge Base ─────────────────────────────────

enterprise_kb_router = APIRouter(prefix="/enterprise/knowledge-base", tags=["enterprise"])


def _enterprise_kb_dir(tenant_id: str) -> Path:
    local_root = settings.STORAGE_LOCAL_ROOT or settings.AGENT_DATA_DIR
    return Path(local_root) / f"enterprise_info_{tenant_id}" / "knowledge_base"


def _enterprise_info_dir(tenant_id: str) -> Path:
    local_root = settings.STORAGE_LOCAL_ROOT or settings.AGENT_DATA_DIR
    return Path(local_root) / f"enterprise_info_{tenant_id}"


def _enterprise_storage_key(tenant_id: str, rel_path: str = "") -> str:
    prefix = f"enterprise_info_{tenant_id}"
    rel = normalize_storage_key(rel_path)
    return f"{prefix}/{rel}" if rel else prefix


@enterprise_kb_router.get("/files")
async def list_enterprise_kb_files(
    path: str = "",
    current_user: User = Depends(get_current_user),
):
    """List files in enterprise knowledge base (tenant-scoped)."""
    if not current_user.tenant_id:
        return []
    storage = get_storage_backend()
    storage_key = _enterprise_storage_key(str(current_user.tenant_id), path)
    if not await storage.exists(storage_key) or not await storage.is_dir(storage_key):
        return []

    items = []
    for entry in await storage.list_dir(storage_key):
        if entry.name == '.gitkeep':
            continue
        rel = str(Path(entry.key).relative_to(f"enterprise_info_{current_user.tenant_id}"))
        items.append({
            "name": entry.name,
            "path": rel,
            "is_dir": entry.is_dir,
            "size": entry.size,
            "url": f"/api/enterprise/knowledge-base/download?path={rel}" if not entry.is_dir else None
        })
    return items


@enterprise_kb_router.post("/upload")
async def upload_enterprise_kb_file(
    file: UploadFileType = FastFile(...),
    sub_path: str = "",
    current_user: User = Depends(get_current_user),
):
    """Upload a file to enterprise knowledge base (tenant-scoped)."""
    # Only admin can upload to enterprise KB
    if current_user.role not in ("platform_admin", "org_admin"):
        raise HTTPException(status_code=403, detail="Only admins can upload to enterprise knowledge base")
    if not current_user.tenant_id:
        raise HTTPException(status_code=400, detail="No tenant associated")

    filename = file.filename or "unnamed"
    filename = filename.replace("/", "_").replace("\\", "_")
    storage = get_storage_backend()
    rel_path = f"{sub_path}/{filename}" if sub_path else filename
    storage_key = _enterprise_storage_key(str(current_user.tenant_id), rel_path)

    content = await file.read()
    await storage.write_bytes(storage_key, content, content_type=guess_content_type(filename))

    # Auto-extract text from non-text files
    extracted_path = None
    from app.services.text_extractor import needs_extraction, save_extracted_text
    if needs_extraction(filename):
        save_path = await ensure_local_path(storage_key)
        txt_file = save_extracted_text(save_path, content, filename)
        if txt_file:
            extracted_path = f"{sub_path}/{txt_file.name}" if sub_path else txt_file.name
            await storage.write_bytes(
                _enterprise_storage_key(str(current_user.tenant_id), extracted_path),
                txt_file.read_bytes(),
                content_type="text/plain; charset=utf-8",
            )
    return {
        "status": "ok",
        "path": rel_path,
        "url": f"/api/enterprise/knowledge-base/download?path={rel_path}",
        "filename": filename,
        "size": len(content),
        "extracted_text_path": extracted_path,
    }


@enterprise_kb_router.get("/content")
async def read_enterprise_file(
    path: str,
    current_user: User = Depends(get_current_user),
):
    """Read content of an enterprise knowledge base file (tenant-scoped)."""
    if not current_user.tenant_id:
        raise HTTPException(status_code=400, detail="No tenant associated")
    storage = get_storage_backend()
    storage_key = _enterprise_storage_key(str(current_user.tenant_id), path)
    if not await storage.exists(storage_key) or not await storage.is_file(storage_key):
        raise HTTPException(status_code=404, detail="File not found")

    content = await _read_text_or_binary_placeholder(storage, storage_key, path)
    return {"path": path, "content": content}


@enterprise_kb_router.put("/content")
async def write_enterprise_file(
    path: str,
    data: FileWrite,
    current_user: User = Depends(get_current_user),
):
    """Write content to an enterprise file (tenant-scoped)."""
    if current_user.role not in ("platform_admin", "org_admin"):
        raise HTTPException(status_code=403, detail="Only admins can edit enterprise knowledge base")
    if not current_user.tenant_id:
        raise HTTPException(status_code=400, detail="No tenant associated")

    storage = get_storage_backend()
    await storage.write_text(_enterprise_storage_key(str(current_user.tenant_id), path), data.content, encoding="utf-8")
    return {"status": "ok", "path": path}


@enterprise_kb_router.delete("/content")
async def delete_enterprise_file(
    path: str,
    current_user: User = Depends(get_current_user),
):
    """Delete an enterprise knowledge base file (tenant-scoped)."""
    if current_user.role not in ("platform_admin", "org_admin"):
        raise HTTPException(status_code=403, detail="Only admins can delete enterprise knowledge base files")
    if not current_user.tenant_id:
        raise HTTPException(status_code=400, detail="No tenant associated")

    storage = get_storage_backend()
    storage_key = _enterprise_storage_key(str(current_user.tenant_id), path)
    storage_exists = await storage.exists(storage_key)
    storage_is_dir = await storage.is_dir(storage_key)
    if not storage_exists and not storage_is_dir:
        raise HTTPException(status_code=404, detail="File not found")
    if storage_is_dir:
        await storage.delete_tree(storage_key)
    else:
        await storage.delete(storage_key)
    return {"status": "ok", "path": path}


# ─── Agent-level ClawHub / URL Skill Import ─────────────────

class ClawhubImportBody(BaseModel):
    slug: str

class UrlImportBody(BaseModel):
    url: str


async def _write_external_skill_snapshot(
    *,
    agent_id: uuid.UUID,
    folder_name: str,
    files: list[dict],
    source: str,
) -> dict:
    """Write a validated, non-overwriting external Skill snapshot."""
    from app.services.skill_workspace import (
        SkillWorkspaceConflictError,
        SkillWorkspaceError,
        import_skill_snapshot_to_agent_workspace,
    )

    try:
        return await import_skill_snapshot_to_agent_workspace(
            storage=get_storage_backend(),
            agent_prefix=_agent_storage_key(agent_id),
            folder_name=folder_name,
            files=files,
            source=source,
        )
    except SkillWorkspaceConflictError as exc:
        raise HTTPException(status.HTTP_409_CONFLICT, str(exc)) from exc
    except SkillWorkspaceError as exc:
        raise HTTPException(status.HTTP_400_BAD_REQUEST, str(exc)) from exc


@router.post("/import-from-clawhub")
async def agent_import_from_clawhub(
    agent_id: uuid.UUID,
    body: ClawhubImportBody,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Import a skill from ClawHub directly into this agent's skills/ workspace."""
    await check_agent_access(db, current_user, agent_id)

    from app.api.skills import (
        _fetch_clawhub_skill_archive, _fetch_clawhub_skill_meta, _get_clawhub_key,
    )

    slug = body.slug
    tenant_id = str(current_user.tenant_id) if current_user.tenant_id else None
    api_key = await _get_clawhub_key(tenant_id)

    # 1. Fetch metadata from ClawHub
    try:
        meta, meta_base = await _fetch_clawhub_skill_meta(slug, api_key=api_key)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(502, f"Failed to connect to ClawHub: {e}")

    skill_info = meta.get("skill", {})
    moderation = meta.get("moderation") or {}
    if moderation.get("isSuspicious"):
        summary = str(
            moderation.get("summary")
            or "ClawHub moderation flagged this package"
        )[:300]
        raise HTTPException(
            status.HTTP_422_UNPROCESSABLE_CONTENT,
            f"Skill installation blocked by ClawHub moderation: {summary}",
        )

    # 2. Fetch files from the ClawHub archive
    files, _ = await _fetch_clawhub_skill_archive(slug, api_key=api_key, preferred_base=meta_base)

    # 3. Write through the configured storage backend without overwriting
    # customized workspace Skills. Importing instructions does not grant Tools.
    imported = await _write_external_skill_snapshot(
        agent_id=agent_id,
        folder_name=slug,
        files=files,
        source=f"clawhub:{slug}",
    )

    return {
        "status": "ok",
        "skill_name": skill_info.get("displayName", slug),
        **imported,
    }


@router.post("/import-from-url")
async def agent_import_from_url(
    agent_id: uuid.UUID,
    body: UrlImportBody,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Import a skill from a GitHub URL directly into this agent's skills/ workspace."""
    await check_agent_access(db, current_user, agent_id)

    from app.api.skills import _parse_github_url, _fetch_github_directory, _get_github_token

    parsed = _parse_github_url(body.url)
    if not parsed:
        raise HTTPException(400, "Invalid GitHub URL")

    owner, repo, branch, path = parsed["owner"], parsed["repo"], parsed["branch"], parsed["path"]
    tenant_id = str(current_user.tenant_id) if current_user.tenant_id else None
    token = await _get_github_token(tenant_id)
    files = await _fetch_github_directory(owner, repo, path, branch, token)
    if not files:
        raise HTTPException(404, "No files found")

    # Derive folder name
    folder_name = path.rstrip("/").split("/")[-1] if path else repo

    imported = await _write_external_skill_snapshot(
        agent_id=agent_id,
        folder_name=folder_name,
        files=files,
        # Persist only the parsed repository identity. The submitted URL may
        # contain credentials or tracking query parameters and must not be
        # copied into an Agent-readable workspace manifest.
        source=f"github:{owner}/{repo}@{branch}:{path or '.'}",
    )

    return {
        "status": "ok",
        **imported,
    }
