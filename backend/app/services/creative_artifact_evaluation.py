"""Machine-observable artifact checks for creative evaluation.

These checks deliberately stop at facts that can be proven from the files.
Semantic accuracy, visual quality, watermark detection, and identity fidelity
remain unknown until an appropriate evaluator supplies evidence.
"""

from __future__ import annotations

from io import BytesIO
import hashlib
from pathlib import Path
from typing import Any, Mapping
import zipfile

from pydantic import BaseModel, ConfigDict, Field

from app.services.creative_evaluation import (
    CreativeQualityEvaluation,
    CreativeScenario,
    CreativeModality,
    HardGateObservation,
    score_quality_evaluation,
)


MAX_EVALUATION_ARTIFACT_BYTES = 200 * 1024 * 1024


class CreativeArtifactContract(BaseModel):
    model_config = ConfigDict(extra="forbid", frozen=True)

    modality: CreativeModality
    aspect_ratio: str
    duration_seconds: float | None = Field(default=None, gt=0)
    page_count: int | None = Field(default=None, gt=0)
    audio_required: bool = False
    reference_identity_required: bool = False
    editable_required: bool = True
    preview_required: bool = True
    minimum_picture_coverage_ratio: float | None = Field(
        default=None,
        ge=0,
        le=1,
    )


class ObservedArtifactFile(BaseModel):
    model_config = ConfigDict(extra="forbid", frozen=True)

    artifact_type: str
    content_sha256: str
    size_bytes: int


class CreativeArtifactObservation(BaseModel):
    model_config = ConfigDict(extra="forbid", frozen=True)

    schema_version: str = "1.0.0"
    modality: CreativeModality
    files: dict[str, ObservedArtifactFile]
    facts: dict[str, Any]
    hard_gates: dict[str, HardGateObservation]
    warnings: tuple[str, ...] = ()


def _read_artifact(path: Path, *, expected_suffixes: set[str]) -> bytes:
    if path.is_symlink() or not path.is_file():
        raise ValueError(f"Artifact is not a regular file: {path.name}")
    if path.suffix.lower() not in expected_suffixes:
        raise ValueError(f"Artifact has unexpected file type: {path.name}")
    size = path.stat().st_size
    if size <= 0 or size > MAX_EVALUATION_ARTIFACT_BYTES:
        raise ValueError(f"Artifact size is outside evaluation limits: {path.name}")
    return path.read_bytes()


def _file_observation(artifact_type: str, data: bytes) -> ObservedArtifactFile:
    return ObservedArtifactFile(
        artifact_type=artifact_type,
        content_sha256=hashlib.sha256(data).hexdigest(),
        size_bytes=len(data),
    )


def _target_ratio(value: str) -> float | None:
    try:
        width, height = value.split(":", maxsplit=1)
        width_value = float(width)
        height_value = float(height)
    except (AttributeError, TypeError, ValueError):
        return None
    if width_value <= 0 or height_value <= 0:
        return None
    return width_value / height_value


def _ratio_matches(width: float, height: float, expected: str, *, tolerance: float) -> bool:
    target = _target_ratio(expected)
    if target is None or width <= 0 or height <= 0:
        return False
    return abs(width / height - target) / target <= tolerance


def _unknown_gate(reason: str) -> HardGateObservation:
    return HardGateObservation(passed=None, evidence=(reason,))


def _observe_image(
    contract: CreativeArtifactContract,
    path: Path,
) -> CreativeArtifactObservation:
    from PIL import Image

    data = _read_artifact(
        path,
        expected_suffixes={".png", ".jpg", ".jpeg", ".webp"},
    )
    try:
        with Image.open(BytesIO(data)) as image:
            image.verify()
        with Image.open(BytesIO(data)) as image:
            width, height = image.size
            image_format = str(image.format or "").upper()
            image_mode = image.mode
    except Exception:
        return CreativeArtifactObservation(
            modality="image",
            files={"image": _file_observation("image", data)},
            facts={},
            hard_gates={
                "artifact_decodable": HardGateObservation(
                    passed=False,
                    evidence=("Pillow could not decode and verify the image",),
                ),
                "aspect_ratio_match": HardGateObservation(passed=None),
                "fact_safety": _unknown_gate("Requires semantic review"),
                "reference_identity_when_required": _unknown_gate(
                    "Requires reference-aware visual review"
                ),
                "no_unrequested_watermark": _unknown_gate(
                    "Requires watermark/OCR review"
                ),
            },
        )

    ratio_match = _ratio_matches(
        width,
        height,
        contract.aspect_ratio,
        tolerance=0.02,
    )
    identity_gate = (
        _unknown_gate("Requires reference-aware visual review")
        if contract.reference_identity_required
        else HardGateObservation(
            passed=True,
            evidence=("No reference identity was required by the contract",),
        )
    )
    return CreativeArtifactObservation(
        modality="image",
        files={"image": _file_observation("image", data)},
        facts={
            "width": width,
            "height": height,
            "format": image_format,
            "mode": image_mode,
            "actual_aspect_ratio": round(width / height, 6),
            "expected_aspect_ratio": contract.aspect_ratio,
        },
        hard_gates={
            "artifact_decodable": HardGateObservation(
                passed=True,
                evidence=(f"Decoded {image_format} {width}x{height}",),
            ),
            "aspect_ratio_match": HardGateObservation(
                passed=ratio_match,
                evidence=(
                    f"Observed {width}:{height}; expected {contract.aspect_ratio}",
                ),
            ),
            "fact_safety": _unknown_gate("Requires semantic review"),
            "reference_identity_when_required": identity_gate,
            "no_unrequested_watermark": _unknown_gate(
                "Requires watermark/OCR review"
            ),
        },
    )


async def _observe_video(
    contract: CreativeArtifactContract,
    path: Path,
) -> CreativeArtifactObservation:
    from app.services.media_assets import (
        validate_generated_video,
        validate_video_delivery_contract,
    )

    data = _read_artifact(path, expected_suffixes={".mp4"})
    file_observation = _file_observation("mp4", data)
    try:
        info = await validate_generated_video(
            data,
            label="Creative evaluation video",
            require_browser_safe=True,
        )
    except (OSError, TypeError, ValueError) as exc:
        return CreativeArtifactObservation(
            modality="video",
            files={"mp4": file_observation},
            facts={},
            hard_gates={
                "artifact_decodable": HardGateObservation(
                    passed=False,
                    evidence=(type(exc).__name__,),
                ),
                "duration_and_aspect_match": HardGateObservation(passed=None),
                "fact_safety": _unknown_gate("Requires semantic review"),
                "audio_contract_match": HardGateObservation(passed=None),
                "no_unrequested_watermark": _unknown_gate(
                    "Requires watermark/frame review"
                ),
            },
        )

    duration_aspect_match = True
    try:
        validate_video_delivery_contract(
            info,
            expected_duration_seconds=contract.duration_seconds,
            expected_aspect_ratio=contract.aspect_ratio,
            require_audio=False,
        )
    except (TypeError, ValueError):
        duration_aspect_match = False
    audio_matches = not contract.audio_required or info.audio_codec_name is not None
    return CreativeArtifactObservation(
        modality="video",
        files={"mp4": file_observation},
        facts={
            "width": info.width,
            "height": info.height,
            "duration_seconds": round(info.duration_seconds, 3),
            "video_codec": info.codec_name,
            "pixel_format": info.pixel_format,
            "audio_codec": info.audio_codec_name,
            "fast_start": info.fast_start,
            "expected_aspect_ratio": contract.aspect_ratio,
            "expected_duration_seconds": contract.duration_seconds,
            "audio_required": contract.audio_required,
        },
        hard_gates={
            "artifact_decodable": HardGateObservation(
                passed=True,
                evidence=(
                    "Browser-safe video stream decoded",
                    f"{info.codec_name}/{info.pixel_format}",
                ),
            ),
            "duration_and_aspect_match": HardGateObservation(
                passed=duration_aspect_match,
                evidence=(
                    f"Observed {info.width}x{info.height}, "
                    f"{info.duration_seconds:.3f}s",
                ),
            ),
            "fact_safety": _unknown_gate("Requires semantic review"),
            "audio_contract_match": HardGateObservation(
                passed=audio_matches,
                evidence=(
                    f"audio_required={contract.audio_required}; "
                    f"audio_codec={info.audio_codec_name or 'none'}",
                ),
            ),
            "no_unrequested_watermark": _unknown_gate(
                "Requires watermark/frame review"
            ),
        },
    )


def _pptx_facts(data: bytes) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    with zipfile.ZipFile(BytesIO(data)) as archive:
        if archive.testzip() is not None:
            raise ValueError("PPTX archive is corrupt")
        names = set(archive.namelist())
        if not {"[Content_Types].xml", "ppt/presentation.xml"} <= names:
            raise ValueError("PPTX package is incomplete")
    presentation = Presentation(BytesIO(data))
    if not presentation.slides:
        raise ValueError("PPTX has no slides")
    width = int(presentation.slide_width or 0)
    height = int(presentation.slide_height or 0)
    if width <= 0 or height <= 0:
        raise ValueError("PPTX slide size is invalid")

    editable_slides = 0
    full_slide_image_only_slides = 0
    text_shape_count = 0
    picture_count = 0
    picture_count_by_slide: list[int] = []
    granular_picture_count_by_slide: list[int] = []
    picture_coverage_ratio_by_slide: list[float] = []
    picture_asset_hashes: set[str] = set()
    full_slide_picture_slides = 0
    slide_area = width * height
    for slide in presentation.slides:
        slide_has_editable_content = False
        slide_has_full_picture = False
        slide_has_other_visible_content = False
        slide_picture_count = 0
        slide_granular_picture_count = 0
        slide_picture_area = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
                slide_picture_count += 1
                try:
                    picture_asset_hashes.add(
                        hashlib.sha256(shape.image.blob).hexdigest()
                    )
                except (AttributeError, KeyError, OSError, ValueError):
                    pass
                shape_left = int(shape.left)
                shape_top = int(shape.top)
                shape_right = shape_left + int(shape.width)
                shape_bottom = shape_top + int(shape.height)
                visible_left = max(
                    0,
                    min(width, min(shape_left, shape_right)),
                )
                visible_top = max(
                    0,
                    min(height, min(shape_top, shape_bottom)),
                )
                visible_right = max(
                    0,
                    min(width, max(shape_left, shape_right)),
                )
                visible_bottom = max(
                    0,
                    min(height, max(shape_top, shape_bottom)),
                )
                if visible_right > visible_left and visible_bottom > visible_top:
                    slide_picture_area += (
                        visible_right - visible_left
                    ) * (visible_bottom - visible_top)
                covers_slide = (
                    int(shape.left) <= width * 0.02
                    and int(shape.top) <= height * 0.02
                    and int(shape.width) >= width * 0.96
                    and int(shape.height) >= height * 0.96
                )
                slide_has_full_picture = slide_has_full_picture or covers_slide
                if not covers_slide:
                    slide_granular_picture_count += 1
                continue
            text = ""
            if getattr(shape, "has_text_frame", False):
                text = str(shape.text or "").strip()
            if text:
                text_shape_count += 1
                slide_has_editable_content = True
            if shape.shape_type in {
                MSO_SHAPE_TYPE.CHART,
                MSO_SHAPE_TYPE.TABLE,
                MSO_SHAPE_TYPE.AUTO_SHAPE,
                MSO_SHAPE_TYPE.GROUP,
            }:
                slide_has_editable_content = True
                slide_has_other_visible_content = True
        if slide_has_editable_content:
            editable_slides += 1
        if slide_has_full_picture and not (
            slide_has_editable_content or slide_has_other_visible_content
        ):
            full_slide_image_only_slides += 1
        if slide_has_full_picture:
            full_slide_picture_slides += 1
        picture_count_by_slide.append(slide_picture_count)
        granular_picture_count_by_slide.append(slide_granular_picture_count)
        picture_coverage_ratio_by_slide.append(
            round(min(slide_picture_area / slide_area, 1.0), 6)
        )
    page_count = len(presentation.slides)
    return {
        "page_count": page_count,
        "width": width,
        "height": height,
        "actual_aspect_ratio": round(width / height, 6),
        "editable_slide_ratio": round(editable_slides / page_count, 6),
        "full_slide_image_only_ratio": round(
            full_slide_image_only_slides / page_count,
            6,
        ),
        "text_shape_count": text_shape_count,
        "picture_count": picture_count,
        "picture_count_by_slide": picture_count_by_slide,
        "granular_picture_count_by_slide": granular_picture_count_by_slide,
        "distinct_picture_asset_count": len(picture_asset_hashes),
        "slides_with_pictures": sum(
            count > 0 for count in picture_count_by_slide
        ),
        "slides_with_pictures_ratio": round(
            sum(count > 0 for count in picture_count_by_slide) / page_count,
            6,
        ),
        "pictures_per_slide_mean": round(picture_count / page_count, 6),
        "pictures_per_slide_min": min(picture_count_by_slide),
        "pictures_per_slide_max": max(picture_count_by_slide),
        "slides_with_granular_pictures": sum(
            count > 0 for count in granular_picture_count_by_slide
        ),
        "slides_with_granular_pictures_ratio": round(
            sum(count > 0 for count in granular_picture_count_by_slide)
            / page_count,
            6,
        ),
        "full_slide_picture_slides": full_slide_picture_slides,
        "full_slide_picture_ratio": round(
            full_slide_picture_slides / page_count,
            6,
        ),
        "picture_coverage_ratio_by_slide": picture_coverage_ratio_by_slide,
        "picture_coverage_ratio_mean": round(
            sum(picture_coverage_ratio_by_slide) / page_count,
            6,
        ),
        "picture_coverage_ratio_min": min(picture_coverage_ratio_by_slide),
        "picture_coverage_ratio_max": max(picture_coverage_ratio_by_slide),
    }


def _pdf_facts(data: bytes) -> dict[str, Any]:
    import fitz

    with fitz.open(stream=data, filetype="pdf") as document:
        if document.page_count <= 0:
            raise ValueError("PDF has no pages")
        page_sizes = [
            (float(page.rect.width), float(page.rect.height)) for page in document
        ]
    width, height = page_sizes[0]
    if width <= 0 or height <= 0:
        raise ValueError("PDF page size is invalid")
    if any(
        abs(other_width - width) > 0.5
        or abs(other_height - height) > 0.5
        for other_width, other_height in page_sizes
    ):
        raise ValueError("PDF page sizes are inconsistent")
    return {
        "page_count": len(page_sizes),
        "width": width,
        "height": height,
        "actual_aspect_ratio": round(width / height, 6),
    }


def _observe_presentation(
    contract: CreativeArtifactContract,
    artifacts: Mapping[str, Path],
) -> CreativeArtifactObservation:
    files: dict[str, ObservedArtifactFile] = {}
    facts: dict[str, Any] = {}
    pptx_valid = False
    pdf_valid = False
    warnings: list[str] = []

    pptx_path = artifacts.get("pptx")
    if pptx_path is not None:
        pptx_data = _read_artifact(pptx_path, expected_suffixes={".pptx"})
        files["pptx"] = _file_observation("pptx", pptx_data)
        try:
            facts["pptx"] = _pptx_facts(pptx_data)
            pptx_valid = True
        except (OSError, TypeError, ValueError, zipfile.BadZipFile) as exc:
            warnings.append(f"pptx_invalid:{type(exc).__name__}")
    else:
        warnings.append("pptx_missing")

    pdf_path = artifacts.get("pdf")
    if pdf_path is not None:
        pdf_data = _read_artifact(pdf_path, expected_suffixes={".pdf"})
        files["pdf"] = _file_observation("pdf", pdf_data)
        try:
            facts["pdf"] = _pdf_facts(pdf_data)
            pdf_valid = True
        except (OSError, TypeError, ValueError) as exc:
            warnings.append(f"pdf_invalid:{type(exc).__name__}")
    elif contract.preview_required:
        warnings.append("pdf_missing")

    preview_contract_valid = pptx_valid and (
        pdf_valid or not contract.preview_required
    )
    page_aspect_match: bool | None = None
    if pptx_valid:
        pptx = facts["pptx"]
        page_aspect_match = _ratio_matches(
            pptx["width"],
            pptx["height"],
            contract.aspect_ratio,
            tolerance=0.015,
        )
        if contract.page_count is not None:
            page_aspect_match = (
                page_aspect_match
                and pptx["page_count"] == contract.page_count
            )
        if pdf_valid:
            pdf = facts["pdf"]
            page_aspect_match = (
                page_aspect_match
                and pdf["page_count"] == pptx["page_count"]
                and _ratio_matches(
                    pdf["width"],
                    pdf["height"],
                    contract.aspect_ratio,
                    tolerance=0.015,
                )
            )

    editable_match: bool | None = None
    if pptx_valid:
        editable_match = (
            not contract.editable_required
            or (
                facts["pptx"]["editable_slide_ratio"] >= 0.8
                and facts["pptx"]["full_slide_image_only_ratio"] <= 0.2
            )
        )

    picture_coverage_match: bool | None = None
    if contract.minimum_picture_coverage_ratio is not None:
        if pptx_valid:
            observed_coverage = float(
                facts["pptx"].get("picture_coverage_ratio_mean") or 0.0
            )
            picture_coverage_match = (
                observed_coverage >= contract.minimum_picture_coverage_ratio
            )
            picture_coverage_evidence = (
                f"observed_mean={observed_coverage:.6f}; "
                f"required_mean={contract.minimum_picture_coverage_ratio:.6f}"
            )
        else:
            picture_coverage_evidence = "PPTX facts unavailable"
    else:
        picture_coverage_evidence = "No minimum picture coverage requested"

    hard_gates: dict[str, HardGateObservation] = {
        "pptx_and_preview_valid": HardGateObservation(
            passed=preview_contract_valid,
            evidence=(
                f"pptx_valid={pptx_valid}; pdf_valid={pdf_valid}; "
                f"preview_required={contract.preview_required}",
            ),
        ),
        "page_count_and_aspect_match": HardGateObservation(
            passed=page_aspect_match,
            evidence=(
                f"expected_page_count={contract.page_count}; "
                f"expected_aspect_ratio={contract.aspect_ratio}",
            ),
        ),
        "fact_safety": _unknown_gate("Requires semantic/source review"),
        "no_text_overflow": _unknown_gate(
            "Requires rendered-slide geometry and visual review"
        ),
        "source_traceability": _unknown_gate(
            "Requires source manifest and citation review"
        ),
        "editability": HardGateObservation(
            passed=editable_match,
            evidence=(
                "Measured editable shapes and full-slide image-only ratio",
            ),
        ),
    }
    if contract.minimum_picture_coverage_ratio is not None:
        hard_gates["minimum_picture_coverage"] = HardGateObservation(
            passed=picture_coverage_match,
            evidence=(picture_coverage_evidence,),
        )

    return CreativeArtifactObservation(
        modality="presentation",
        files=files,
        facts=facts,
        hard_gates=hard_gates,
        warnings=tuple(warnings),
    )


async def observe_creative_artifacts(
    contract: CreativeArtifactContract,
    artifacts: Mapping[str, Path],
) -> CreativeArtifactObservation:
    """Inspect a provider-neutral artifact set without fabricating semantic passes."""

    if contract.modality == "image":
        path = artifacts.get("image")
        if path is None:
            raise ValueError("Image evaluation requires an 'image' artifact")
        return _observe_image(contract, path)
    if contract.modality == "video":
        path = artifacts.get("mp4")
        if path is None:
            raise ValueError("Video evaluation requires an 'mp4' artifact")
        return await _observe_video(contract, path)
    return _observe_presentation(contract, artifacts)


def score_artifact_observation(
    scenario: CreativeScenario,
    observation: CreativeArtifactObservation,
) -> CreativeQualityEvaluation:
    """Feed proven hard gates into the common scorer.

    Machine observations intentionally do not populate subjective dimensions,
    so a structurally valid artifact remains ``incomplete`` until reviewed.
    """

    if scenario.modality != observation.modality:
        raise ValueError("Scenario and artifact observation modalities differ")
    scoring_scenario = scenario
    if (
        "minimum_picture_coverage" in observation.hard_gates
        and "minimum_picture_coverage" not in scenario.hard_gates
    ):
        # The threshold is an optional artifact contract rather than a
        # requirement of every presentation scenario.  When a caller opts
        # into it, include the observed gate in scoring instead of silently
        # dropping an enforceable structural failure.
        scoring_scenario = scenario.model_copy(
            update={
                "hard_gates": (
                    *scenario.hard_gates,
                    "minimum_picture_coverage",
                )
            }
        )
    return score_quality_evaluation(
        scoring_scenario,
        hard_gates=observation.hard_gates,
        dimensions={},
    )


__all__ = [
    "CreativeArtifactContract",
    "CreativeArtifactObservation",
    "ObservedArtifactFile",
    "observe_creative_artifacts",
    "score_artifact_observation",
]
